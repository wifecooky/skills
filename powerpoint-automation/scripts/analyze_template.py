# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Analyze PowerPoint template and generate layout mapping configuration.

This script inspects the slide master layouts in a template and creates
a JSON configuration file that maps slide types to appropriate layouts.

Usage:
    python scripts/analyze_template.py <template.pptx> [output.json]

Examples:
    python scripts/analyze_template.py templates/sample-ppf.pptx
    python scripts/analyze_template.py templates/sample-ppf.pptx output_manifest/sample-ppf_layouts.json
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import sys
from pathlib import Path


def analyze_layout(layout, idx: int) -> dict:
    """Analyze a single slide layout and return its characteristics."""
    name = layout.name or f'Layout_{idx}'
    
    # Collect placeholder info
    placeholders = []
    has_title = False
    has_subtitle = False
    has_body = False
    has_content = False
    has_picture = False
    body_count = 0
    
    for shape in layout.placeholders:
        ph_type = str(shape.placeholder_format.type)
        ph_idx = shape.placeholder_format.idx
        
        ph_info = {
            'idx': ph_idx,
            'type': ph_type.replace('PLACEHOLDER_TYPE.', ''),
        }
        placeholders.append(ph_info)
        
        if 'TITLE' in ph_type or 'CENTER_TITLE' in ph_type:
            has_title = True
        if 'SUBTITLE' in ph_type:
            has_subtitle = True
        if 'BODY' in ph_type:
            has_body = True
            body_count += 1
        if 'OBJECT' in ph_type or 'CONTENT' in ph_type:
            has_content = True
        if 'PICTURE' in ph_type:
            has_picture = True
    
    # Determine layout category - order matters! More specific patterns first
    category = 'unknown'
    name_lower = name.lower()
    
    # Explicit closing patterns (most specific first)
    if 'closing' in name_lower:
        category = 'closing'
    # Title slide patterns (before general title checks)
    elif 'title slide' in name_lower or 'タイトル スライド' in name_lower or 'タイトルスライド' in name_lower:
        category = 'title'
    # Section patterns
    elif 'section' in name_lower or 'セクション' in name_lower or 'divider' in name_lower:
        category = 'section'
    # Agenda patterns
    elif 'agenda' in name_lower or 'アジェンダ' in name_lower:
        category = 'agenda'
    # Content patterns
    elif 'title and content' in name_lower or 'タイトルとコンテンツ' in name_lower:
        category = 'content'
    # Two column
    elif 'two column' in name_lower or '2列' in name_lower or '2 column' in name_lower:
        category = 'two_column'
    # Three column
    elif 'three column' in name_lower or '3列' in name_lower or '3 column' in name_lower:
        category = 'three_column'
    # Code/developer layouts
    elif 'code' in name_lower or 'developer' in name_lower:
        category = 'code'
    # Quote layouts
    elif 'quote' in name_lower:
        category = 'quote'
    # Photo/image layouts
    elif 'photo' in name_lower or 'picture' in name_lower or 'image' in name_lower or '50/50' in name_lower:
        category = 'photo'
    elif 'blank' in name_lower or '白紙' in name_lower:
        category = 'blank'
    elif 'title only' in name_lower:
        category = 'title_only'
    elif has_title and (has_body or has_content):
        category = 'content'
    elif has_title and has_subtitle and not has_body:
        category = 'title'
    
    return {
        'index': idx,
        'name': name,
        'category': category,
        'has_title': has_title,
        'has_subtitle': has_subtitle,
        'has_body': has_body,
        'has_content': has_content,
        'has_picture': has_picture,
        'body_count': body_count,
        'placeholders': placeholders,
    }


def generate_layout_mapping(layouts_info: list) -> dict:
    """Generate recommended layout mapping based on analysis."""
    mapping = {
        'title': None,
        'content': None,
        'section': None,
        'agenda': None,
        'summary': None,
        'closing': None,
        'two_column': None,
        'code': None,
        'photo': None,
        'blank': None,
        'title_only': None,  # Added: for slides with title but no content
    }
    
    # Find best match for each category
    for layout in layouts_info:
        cat = layout['category']
        idx = layout['index']
        
        # Prefer first match for each category
        if cat == 'title' and mapping['title'] is None:
            mapping['title'] = idx
        elif cat == 'content' and mapping['content'] is None:
            mapping['content'] = idx
        elif cat == 'section' and mapping['section'] is None:
            mapping['section'] = idx
        elif cat == 'agenda' and mapping['agenda'] is None:
            mapping['agenda'] = idx
        elif cat == 'closing' and mapping['closing'] is None:
            mapping['closing'] = idx
        elif cat == 'two_column' and mapping['two_column'] is None:
            mapping['two_column'] = idx
        elif cat == 'code' and mapping['code'] is None:
            mapping['code'] = idx
        elif cat == 'photo' and mapping['photo'] is None:
            mapping['photo'] = idx
        elif cat == 'blank' and mapping['blank'] is None:
            mapping['blank'] = idx
        elif cat == 'title_only' and mapping['title_only'] is None:
            mapping['title_only'] = idx
    
    # Set fallbacks
    if mapping['content'] is None:
        mapping['content'] = 1  # Default fallback
    if mapping['title'] is None:
        mapping['title'] = 0
    if mapping['summary'] is None:
        mapping['summary'] = mapping['content']
    if mapping['section'] is None:
        mapping['section'] = mapping['content']
    if mapping['agenda'] is None:
        mapping['agenda'] = mapping['content']
    if mapping['closing'] is None:
        mapping['closing'] = mapping['section']
    if mapping['two_column'] is None:
        mapping['two_column'] = mapping['content']
    if mapping['code'] is None:
        mapping['code'] = mapping['content']
    if mapping['photo'] is None:
        mapping['photo'] = mapping['content']
    if mapping['blank'] is None:
        mapping['blank'] = mapping['content']
    if mapping['title_only'] is None:
        # Fallback: prefer blank, then section, then content
        mapping['title_only'] = mapping['blank'] or mapping['section'] or mapping['content']
    
    return mapping


def analyze_template(template_path: str) -> dict:
    """Analyze a PowerPoint template and return full analysis."""
    prs = Presentation(template_path)
    
    result = {
        'template': str(Path(template_path).name),
        'template_path': str(template_path),
        'slide_width_inches': round(prs.slide_width.inches, 2),
        'slide_height_inches': round(prs.slide_height.inches, 2),
        'aspect_ratio': '16:9' if abs(prs.slide_width.inches / prs.slide_height.inches - 16/9) < 0.1 else '4:3',
        'masters': [],
        'layouts': [],
        'layout_mapping': {},
    }
    
    all_layouts = []
    
    for master_idx, master in enumerate(prs.slide_masters):
        master_info = {
            'index': master_idx,
            'name': master.name or f'Master_{master_idx}',
            'layout_count': len(master.slide_layouts),
        }
        result['masters'].append(master_info)
    
    # Analyze all layouts (from first master)
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = analyze_layout(layout, idx)
        all_layouts.append(layout_info)
    
    result['layouts'] = all_layouts
    result['layout_mapping'] = generate_layout_mapping(all_layouts)
    
    return result


def print_analysis(analysis: dict) -> None:
    """Print analysis in human-readable format."""
    print(f"\n{'='*60}")
    print(f"Template: {analysis['template']}")
    print(f"Size: {analysis['slide_width_inches']}\" x {analysis['slide_height_inches']}\" ({analysis['aspect_ratio']})")
    print(f"{'='*60}\n")
    
    print("📋 Recommended Layout Mapping:")
    print("-" * 40)
    for slide_type, idx in analysis['layout_mapping'].items():
        if idx is not None:
            layout_name = analysis['layouts'][idx]['name'] if idx < len(analysis['layouts']) else '?'
            print(f"  {slide_type:15} → [{idx:2}] {layout_name}")
    print()
    
    print("📊 All Layouts by Category:")
    print("-" * 40)
    
    # Group by category
    by_category = {}
    for layout in analysis['layouts']:
        cat = layout['category']
        if cat not in by_category:
            by_category[cat] = []
        by_category[cat].append(layout)
    
    for cat in sorted(by_category.keys()):
        print(f"\n  [{cat}]")
        for layout in by_category[cat][:5]:  # Show max 5 per category
            flags = []
            if layout['has_title']:
                flags.append('T')
            if layout['has_body']:
                flags.append(f"B{layout['body_count']}")
            if layout['has_content']:
                flags.append('C')
            if layout['has_picture']:
                flags.append('P')
            flag_str = ','.join(flags) if flags else '-'
            print(f"    [{layout['index']:2}] {layout['name'][:40]:40} ({flag_str})")
        if len(by_category[cat]) > 5:
            print(f"    ... and {len(by_category[cat]) - 5} more")
    
    print()


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    
    template_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    # Auto-generate output path if not specified
    if output_path is None:
        template_stem = Path(template_path).stem
        output_path = f"output_manifest/{template_stem}_layouts.json"
    
    # Analyze
    print(f"Analyzing: {template_path}")
    analysis = analyze_template(template_path)
    
    # Print summary
    print_analysis(analysis)
    
    # Save JSON
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(analysis, f, ensure_ascii=False, indent=2)
    
    print(f"✅ Saved: {output_path}")
    print(f"\nTo use this mapping, either:")
    print(f"  1. Edit the JSON and pass to create_from_template.py --config {output_path}")
    print(f"  2. Copy layout_mapping values to your content JSON")


if __name__ == '__main__':
    main()
