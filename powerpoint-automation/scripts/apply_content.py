#!/usr/bin/env python3
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Apply text replacements and slide management to PowerPoint presentation.

This module provides functionality to:
- Replace text content in PowerPoint presentations based on JSON configuration
- Delete slides (via slides_to_delete or slides_to_keep)
- Preserve formatting and support bullet lists, fonts, colors, paragraph styling

Usage:
    python apply_content.py input.pptx replacements.json output.pptx

JSON Structure:
    {
      "slides_to_keep": [0, 1, 5, 10],      // Optional: keep only these slides
      "slides_to_delete": [2, 3, 4],        // Optional: delete these slides
      "slide-0": {                          // Text replacements
        "shape-0": {
          "paragraphs": [{"text": "New text"}]
        }
      }
    }

Note: slides_to_keep takes precedence over slides_to_delete if both are specified.

IMPORTANT:
    This script must use the same shape extraction and sorting logic as
    extract_shapes.py to ensure shape-id consistency. If you modify the shape
    extraction logic in either file, update the other file as well.

Author: aktsmm
License: CC BY-NC 4.0
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt


def get_luminance(r: int, g: int, b: int) -> float:
    """Calculate relative luminance of a color (WCAG formula).
    
    Args:
        r, g, b: RGB values (0-255)
        
    Returns:
        Relative luminance (0.0 to 1.0)
    """
    def channel_luminance(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    
    return 0.2126 * channel_luminance(r) + 0.7152 * channel_luminance(g) + 0.0722 * channel_luminance(b)


def is_dark_background(r: int, g: int, b: int, threshold: float = 0.5) -> bool:
    """Determine if a background color is dark.
    
    Args:
        r, g, b: RGB values (0-255)
        threshold: Luminance threshold (default 0.5)
        
    Returns:
        True if background is dark, False if light
    """
    return get_luminance(r, g, b) < threshold


def get_contrast_color(bg_r: int, bg_g: int, bg_b: int) -> Tuple[int, int, int]:
    """Get optimal text color for given background.
    
    Args:
        bg_r, bg_g, bg_b: Background RGB values (0-255)
        
    Returns:
        Tuple of (r, g, b) for optimal text color
    """
    if is_dark_background(bg_r, bg_g, bg_b):
        return (255, 255, 255)  # White text for dark backgrounds
    else:
        return (51, 51, 51)  # Dark gray (#333333) for light backgrounds


def get_slide_background_color(slide) -> Optional[Tuple[int, int, int]]:
    """Extract background color from a slide.
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        Tuple of (r, g, b) or None if cannot determine
    """
    try:
        # Try to get background fill
        background = slide.background
        fill = background.fill
        
        if fill.type is not None:
            # Solid fill
            if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                color = fill.fore_color
                if color.rgb:
                    return (color.rgb[0], color.rgb[1], color.rgb[2])
        
        # Try slide layout background
        if hasattr(slide, 'slide_layout') and slide.slide_layout:
            layout_bg = slide.slide_layout.background
            if layout_bg and layout_bg.fill:
                fill = layout_bg.fill
                if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                    color = fill.fore_color
                    if color.rgb:
                        return (color.rgb[0], color.rgb[1], color.rgb[2])
        
        # Try slide master background
        if hasattr(slide, 'slide_layout') and slide.slide_layout:
            master = slide.slide_layout.slide_master
            if master and master.background and master.background.fill:
                fill = master.background.fill
                if hasattr(fill, 'fore_color') and fill.fore_color.type is not None:
                    color = fill.fore_color
                    if color.rgb:
                        return (color.rgb[0], color.rgb[1], color.rgb[2])
        
    except Exception:
        pass
    
    return None


def detect_background_from_existing_text(slide) -> Optional[Tuple[int, int, int]]:
    """Detect likely background color from existing text colors.
    
    If most text is white, background is probably dark.
    If most text is dark, background is probably light.
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        Estimated background RGB or None
    """
    white_count = 0
    dark_count = 0
    
    try:
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                continue
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        lum = get_luminance(rgb[0], rgb[1], rgb[2])
                        if lum > 0.8:  # White-ish text
                            white_count += 1
                        elif lum < 0.3:  # Dark text
                            dark_count += 1
    except Exception:
        pass
    
    if white_count > dark_count and white_count > 0:
        # Most text is white -> dark background
        return (30, 30, 60)  # Dark blue-ish (common in presentations)
    elif dark_count > white_count and dark_count > 0:
        # Most text is dark -> light background
        return (240, 240, 245)  # Light gray-ish
    
    return None


def get_optimal_text_color_for_slide(slide) -> Tuple[int, int, int]:
    """Get optimal text color for a slide based on its background.
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        Tuple of (r, g, b) for optimal text color
    """
    # First try to get actual background color
    bg_color = get_slide_background_color(slide)
    
    # If not found, try to infer from existing text colors
    if bg_color is None:
        bg_color = detect_background_from_existing_text(slide)
    
    # If still not found, default to assuming dark background
    # (common in corporate presentations like Microsoft Ignite)
    if bg_color is None:
        return (51, 51, 51)  # Default to dark text (safe for most backgrounds)
    
    return get_contrast_color(*bg_color)


def parse_arguments() -> argparse.Namespace:
    """Parse command line arguments.
    
    Returns:
        Parsed command line arguments.
    """
    parser = argparse.ArgumentParser(
        description="Apply text replacements and slide management to PowerPoint presentation.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python apply_content.py template.pptx replacements.json output.pptx
    Applies text replacements from JSON to create output presentation

=== SLIDE CREATION MODES ===

This script is part of a multi-mode PPTX generation system:

  1. EXISTING PPTX (this script)
     - Use when: Localizing/translating existing presentations
     - Input: Original .pptx + replacements.json
     - Preserves: All original styles, layouts, animations
     - Command: python apply_content.py input.pptx replacements.json output.pptx

  2. TEMPLATE-BASED (extract_shapes.py + this script)
     - Use when: Creating content with corporate templates
     - Input: Template .pptx + inventory.json + replacements.json
     - First: python extract_shapes.py template.pptx inventory.json
     - Then: python apply_content.py template.pptx replacements.json output.pptx

  3. AUTO-SELECT (Orchestrator Agent)
     - Use when: "おまかせ" mode - let AI choose best approach
     - The agent analyzes content and selects optimal method

  4. HTML-BASED (convert_html.js / convert_html_multi.js)
     - Use when: Custom designs, code blocks, complex layouts
     - Input: HTML files in slides/ directory
     - Command: node convert_html_multi.js slides/ output.pptx
     - Offers: Full design freedom, syntax highlighting

=== JSON STRUCTURE ===

{
  "slides_to_keep": [0, 1, 5, 10],      // Optional: keep only these slides (0-indexed)
  "slides_to_delete": [2, 3, 4],        // Optional: delete these slides (0-indexed)
  "slide-0": {
    "shape-0": {
      "paragraphs": [
        {"text": "Title text", "font_size": 32.0},
        {"text": "Bullet item", "bullet": true, "level": 0}
      ]
    }
  }
}

Note: 
  - slides_to_keep takes precedence over slides_to_delete
  - Text replacements are applied BEFORE slide deletion
  - Slide indices refer to the ORIGINAL slide positions
        """,
    )
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("replacements", help="JSON file with replacement data")
    parser.add_argument("output", help="Output PowerPoint file (.pptx)")
    parser.add_argument("--no-auto-shrink", action="store_true",
                        help="Disable automatic font size shrinking for overflow prevention")
    return parser.parse_args()


def load_replacements(json_path: Path) -> Dict[str, Any]:
    """Load replacement data from JSON file.
    
    Args:
        json_path: Path to the JSON file.
        
    Returns:
        Dictionary containing replacement data.
    """
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


def clear_paragraph_formatting(paragraph) -> None:
    """Remove bullet formatting from a paragraph.
    
    Args:
        paragraph: The PowerPoint paragraph object.
    """
    pPr = paragraph._element.get_or_add_pPr()
    
    # Remove existing bullet elements
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ["buChar", "buNone", "buAutoNum", "buFont", "buSzPct", "buSzPts"]:
        for child in list(pPr):
            if child.tag == f"{ns_a}{tag}":
                pPr.remove(child)


def create_bullet_char(pPr, char: str = "•") -> None:
    """Add bullet character formatting to paragraph.
    
    Args:
        pPr: The paragraph properties element.
        char: The bullet character to use.
    """
    from lxml import etree
    
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    buChar = etree.SubElement(pPr, f"{{{ns_a}}}buChar")
    buChar.set("char", char)


def create_no_bullet(pPr) -> None:
    """Add no-bullet formatting to paragraph.
    
    Args:
        pPr: The paragraph properties element.
    """
    from lxml import etree
    
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    buNone = etree.SubElement(pPr, f"{{{ns_a}}}buNone")


def apply_paragraph_formatting(paragraph, para_data: Dict[str, Any]) -> None:
    """Apply formatting properties to a paragraph.
    
    Args:
        paragraph: The PowerPoint paragraph object.
        para_data: Dictionary containing paragraph formatting data.
    """
    text = para_data.get("text", "")
    # Clean text: remove vertical tabs (soft line breaks) that cause "_x000B_" display
    text = text.replace('\x0b', ' ').replace('\u000b', ' ')
    
    # Clear existing formatting
    clear_paragraph_formatting(paragraph)
    pPr = paragraph._element.get_or_add_pPr()
    
    # Handle bullet formatting
    if para_data.get("bullet", False):
        level = para_data.get("level", 0)
        paragraph.level = level
        
        # Calculate indentation based on font size
        font_size = para_data.get("font_size", 18.0)
        level_indent_emu = int((font_size * (1.6 + level * 1.6)) * 12700)
        hanging_indent_emu = int(-font_size * 0.8 * 12700)
        
        # Set indentation
        pPr.attrib["marL"] = str(level_indent_emu)
        pPr.attrib["indent"] = str(hanging_indent_emu)
        
        # Add bullet character
        create_bullet_char(pPr)
        
        # Default to left alignment for bullets
        if "alignment" not in para_data:
            paragraph.alignment = PP_ALIGN.LEFT
    else:
        # Remove indentation for non-bullet text
        pPr.attrib["marL"] = "0"
        pPr.attrib["indent"] = "0"
        create_no_bullet(pPr)
    
    # Apply alignment
    if "alignment" in para_data:
        alignment_map = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFY": PP_ALIGN.JUSTIFY,
        }
        if para_data["alignment"] in alignment_map:
            paragraph.alignment = alignment_map[para_data["alignment"]]
    
    # Apply spacing
    if "space_before" in para_data:
        paragraph.space_before = Pt(para_data["space_before"])
    if "space_after" in para_data:
        paragraph.space_after = Pt(para_data["space_after"])
    
    # Set text on run
    if paragraph.runs:
        run = paragraph.runs[0]
        run.text = text
    else:
        run = paragraph.add_run()
        run.text = text
    
    # Apply font properties
    apply_font_properties(run, para_data)


def apply_font_properties(run, para_data: Dict[str, Any]) -> None:
    """Apply font properties to a text run.
    
    Args:
        run: The PowerPoint run object.
        para_data: Dictionary containing font property data.
    """
    font = run.font
    
    if "font_size" in para_data:
        font.size = Pt(para_data["font_size"])
    if "font_name" in para_data:
        font.name = para_data["font_name"]
    if "bold" in para_data:
        font.bold = para_data["bold"]
    if "italic" in para_data:
        font.italic = para_data["italic"]
    if "underline" in para_data:
        font.underline = para_data["underline"]
    
    # Apply color - prefer RGB, fall back to theme_color
    if "color" in para_data:
        color_hex = para_data["color"].lstrip("#")
        if len(color_hex) == 6:
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            font.color.rgb = RGBColor(r, g, b)
    elif "theme_color" in para_data:
        theme_name = para_data["theme_color"]
        try:
            font.color.theme_color = getattr(MSO_THEME_COLOR, theme_name)
        except AttributeError:
            print(f"  WARNING: Unknown theme color '{theme_name}'")


def emu_to_inches(emu: int) -> float:
    """Convert EMUs (English Metric Units) to inches.
    
    Args:
        emu: Value in EMUs.
        
    Returns:
        Value in inches.
    """
    return emu / 914400


def estimate_text_width(text: str, font_size: float) -> float:
    """Estimate text width in inches based on character count and font size.
    
    Japanese characters are wider than English characters.
    
    Args:
        text: The text to measure.
        font_size: Font size in points.
        
    Returns:
        Estimated width in inches.
    """
    import re
    
    # Count Japanese characters (Hiragana, Katakana, Kanji)
    japanese_pattern = r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]'
    ja_chars = len(re.findall(japanese_pattern, text))
    
    # Count other characters (including spaces, numbers, etc.)
    other_chars = len(text) - ja_chars
    
    # Approximate character widths relative to font size
    # Japanese: ~1.0 em width, English: ~0.5 em width
    # 1 em ≈ font_size in points, 72 points = 1 inch
    ja_width = ja_chars * (font_size / 72) * 0.9  # Japanese chars ~90% of em
    other_width = other_chars * (font_size / 72) * 0.5  # English chars ~50% of em
    
    return ja_width + other_width


def calculate_optimal_font_size(text: str, shape_width_inches: float, 
                                 original_font_size: float,
                                 min_font_size: float = 12.0) -> float:
    """Calculate optimal font size to fit text within shape width.
    
    Args:
        text: The text to fit.
        shape_width_inches: Available width in inches.
        original_font_size: Original font size in points.
        min_font_size: Minimum allowed font size.
        
    Returns:
        Optimal font size in points.
    """
    # Handle multi-line text - check longest line
    lines = text.split('\n')
    max_line = max(lines, key=len) if lines else text
    
    # Estimate width at original font size
    estimated_width = estimate_text_width(max_line, original_font_size)
    
    # If it fits, return original size
    if estimated_width <= shape_width_inches * 0.95:  # 5% margin
        return original_font_size
    
    # Calculate reduction ratio needed
    ratio = (shape_width_inches * 0.95) / estimated_width
    new_font_size = original_font_size * ratio
    
    # Enforce minimum font size
    return max(new_font_size, min_font_size)


def check_overflow_warning(text: str, shape_width_inches: float, 
                           font_size: float, slide_key: str, shape_key: str) -> Optional[str]:
    """Check if text might overflow and return warning message.
    
    Args:
        text: The text to check.
        shape_width_inches: Shape width in inches.
        font_size: Font size in points.
        slide_key: Slide identifier for the warning message.
        shape_key: Shape identifier for the warning message.
        
    Returns:
        Warning message if overflow detected, None otherwise.
    """
    lines = text.split('\n')
    max_line = max(lines, key=len) if lines else text
    estimated_width = estimate_text_width(max_line, font_size)
    
    if estimated_width > shape_width_inches * 1.0:
        return (f"  ⚠️ Overflow risk: {slide_key}.{shape_key} "
                f"(text: {len(max_line)} chars, estimated: {estimated_width:.1f}in > width: {shape_width_inches:.1f}in)")
    
    return None


def extract_shapes_from_slide(slide) -> List[Any]:
    """Extract all shapes from a slide, handling groups recursively.
    
    Uses the same logic as extract_shapes.py to ensure shape-id consistency.
    
    Args:
        slide: The PowerPoint slide object.
        
    Returns:
        List of shapes with text frames, sorted by visual position.
    """
    def get_shapes_recursive(shapes, offset_left: int = 0, offset_top: int = 0):
        result = []
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                group_left = (shape.left or 0) + offset_left
                group_top = (shape.top or 0) + offset_top
                result.extend(get_shapes_recursive(
                    shape.shapes, 
                    group_left,
                    group_top
                ))
            elif shape.has_text_frame:
                # Skip empty text frames
                if not shape.text_frame.text.strip():
                    continue
                    
                # Skip slide numbers and footers
                if shape.is_placeholder:
                    ph_type = str(shape.placeholder_format.type) if shape.placeholder_format.type else ""
                    if "SLIDE_NUMBER" in ph_type or "FOOTER" in ph_type or "DATE" in ph_type:
                        continue
                
                abs_left = (shape.left or 0) + offset_left
                abs_top = (shape.top or 0) + offset_top
                result.append((shape, abs_left, abs_top))
        return result
    
    shapes_with_pos = get_shapes_recursive(slide.shapes)
    
    # Sort by visual position (top to bottom, left to right) - SAME as extract_shapes.py
    shapes_with_pos.sort(key=lambda x: (emu_to_inches(x[2]), emu_to_inches(x[1])))
    
    # Return just the shapes
    return [shape for shape, _, _ in shapes_with_pos]


def apply_replacements(prs: Presentation, replacements: Dict[str, Any], 
                       auto_shrink: bool = True) -> Tuple[int, List[str]]:
    """Apply all replacements to the presentation.
    
    Args:
        prs: The Presentation object.
        replacements: Dictionary of replacements by slide and shape.
        auto_shrink: Whether to automatically shrink font to prevent overflow.
        
    Returns:
        Tuple of (number of shapes modified, list of warnings).
    """
    modified_count = 0
    all_warnings = []
    
    for slide_key, shapes_data in replacements.items():
        # Parse slide index from key (e.g., "slide-0" -> 0)
        if not slide_key.startswith("slide-"):
            continue
        
        try:
            slide_idx = int(slide_key.split("-")[1])
        except (IndexError, ValueError):
            print(f"  WARNING: Invalid slide key '{slide_key}'")
            continue
        
        if slide_idx >= len(prs.slides):
            print(f"  WARNING: Slide index {slide_idx} out of range")
            continue
        
        slide = prs.slides[slide_idx]
        
        # Get shapes in position order
        shapes = extract_shapes_from_slide(slide)
        
        for shape_key, shape_data in shapes_data.items():
            # Parse shape index from key (e.g., "shape-0" -> 0)
            if not shape_key.startswith("shape-"):
                continue
            
            try:
                shape_idx = int(shape_key.split("-")[1])
            except (IndexError, ValueError):
                print(f"  WARNING: Invalid shape key '{shape_key}'")
                continue
            
            if shape_idx >= len(shapes):
                print(f"  WARNING: Shape index {shape_idx} out of range on {slide_key}")
                continue
            
            shape = shapes[shape_idx]
            
            # Apply paragraph replacements
            if "paragraphs" in shape_data:
                warnings = apply_shape_text(
                    shape, shape_data["paragraphs"],
                    slide_key=slide_key, shape_key=shape_key,
                    auto_shrink=auto_shrink
                )
                all_warnings.extend(warnings)
                modified_count += 1
    
    return modified_count, all_warnings


def apply_shape_text(shape, paragraphs: List[Dict[str, Any]], 
                     slide_key: str = "", shape_key: str = "",
                     auto_shrink: bool = True) -> List[str]:
    """Replace all text in a shape with new paragraphs.
    
    Args:
        shape: The PowerPoint shape object.
        paragraphs: List of paragraph data dictionaries.
        slide_key: Slide identifier for logging.
        shape_key: Shape identifier for logging.
        auto_shrink: Whether to automatically shrink font to prevent overflow.
        
    Returns:
        List of warning messages (empty if no warnings).
    """
    warnings = []
    text_frame = shape.text_frame
    
    # Get shape width for overflow detection
    shape_width_inches = emu_to_inches(shape.width) if shape.width else 10.0
    
    # Clear existing paragraphs
    for para in list(text_frame.paragraphs):
        p = para._element
        p.getparent().remove(p)
    
    # Add new paragraphs
    for i, para_data in enumerate(paragraphs):
        if i == 0:
            # First paragraph already exists (empty)
            para = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            para = text_frame.add_paragraph()
        
        # Check for overflow and auto-adjust font size if needed
        text = para_data.get("text", "")
        original_font_size = para_data.get("font_size", 18.0)
        
        if auto_shrink and text:
            optimal_font_size = calculate_optimal_font_size(
                text, shape_width_inches, original_font_size
            )
            
            if optimal_font_size < original_font_size:
                warnings.append(
                    f"  📏 Auto-shrink: {slide_key}.{shape_key} para[{i}] "
                    f"({original_font_size:.0f}pt → {optimal_font_size:.0f}pt)"
                )
                # Update para_data with new font size
                para_data = para_data.copy()
                para_data["font_size"] = optimal_font_size
            else:
                # Still check for overflow warning even if not shrinking
                warning = check_overflow_warning(
                    text, shape_width_inches, original_font_size, slide_key, shape_key
                )
                if warning:
                    warnings.append(warning)
        
        apply_paragraph_formatting(para, para_data)
    
    return warnings


def delete_slides(prs: Presentation, data: Dict[str, Any]) -> int:
    """Delete slides from presentation based on configuration.
    
    Supports two modes:
    - slides_to_keep: Keep only the specified slides (delete all others)
    - slides_to_delete: Delete only the specified slides
    
    slides_to_keep takes precedence if both are specified.
    
    Args:
        prs: The Presentation object.
        data: Dictionary containing slides_to_keep or slides_to_delete.
        
    Returns:
        Number of slides deleted.
    """
    total_slides = len(prs.slides)
    
    # Determine which slides to delete
    slides_to_delete = set()
    
    if "slides_to_keep" in data:
        # Keep only specified slides, delete all others
        slides_to_keep = set(data["slides_to_keep"])
        slides_to_delete = set(range(total_slides)) - slides_to_keep
        print(f"  Mode: Keep {len(slides_to_keep)} slides, delete {len(slides_to_delete)} slides")
    elif "slides_to_delete" in data:
        # Delete specified slides
        slides_to_delete = set(data["slides_to_delete"])
        print(f"  Mode: Delete {len(slides_to_delete)} specified slides")
    else:
        return 0
    
    # Validate slide indices
    valid_indices = set(range(total_slides))
    invalid_indices = slides_to_delete - valid_indices
    if invalid_indices:
        print(f"  WARNING: Skipping invalid slide indices: {sorted(invalid_indices)}")
        slides_to_delete = slides_to_delete & valid_indices
    
    if not slides_to_delete:
        return 0
    
    # Delete slides in reverse order to maintain indices
    # python-pptx doesn't have direct slide delete, use XML manipulation
    deleted_count = 0
    for slide_idx in sorted(slides_to_delete, reverse=True):
        try:
            # Get the slide's rId by finding the relationship
            slide = prs.slides[slide_idx]
            slide_part = slide.part
            
            # Find rId by iterating relationships
            rId = None
            for rel in prs.part.rels.values():
                if rel.target_part == slide_part:
                    rId = rel.rId
                    break
            
            if rId is None:
                print(f"  WARNING: Could not find rId for slide {slide_idx}")
                continue
            
            # Remove from presentation.xml sldIdLst first
            sldIdLst = prs.part._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
            )
            if sldIdLst is not None:
                for sldId in list(sldIdLst):
                    if sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id") == rId:
                        sldIdLst.remove(sldId)
                        break
            
            # Remove the relationship
            prs.part.drop_rel(rId)
            
            deleted_count += 1
        except Exception as e:
            print(f"  WARNING: Failed to delete slide {slide_idx}: {e}")
    
    return deleted_count


def add_summary_slide(prs: Presentation, summary_data: Dict[str, Any], insert_position: int = 1) -> bool:
    """Add a summary/agenda slide at the specified position.
    
    Args:
        prs: The Presentation object.
        summary_data: Dictionary containing 'title', 'items' list, and optional 'color'.
                      If 'color' is 'auto', automatically detect optimal color.
        insert_position: Position to insert the slide (1-indexed, default: after title slide).
        
    Returns:
        True if successful, False otherwise.
    """
    if not summary_data:
        return False
    
    title = summary_data.get("title", "アジェンダ")
    items = summary_data.get("items", [])
    text_color = summary_data.get("color", "auto")  # Default to auto-detect
    
    if not items:
        return False
    
    try:
        # Use Title and Content layout (index 1) or first available layout with content placeholder
        slide_layout = None
        for layout in prs.slide_layouts:
            # Look for layout with title and body placeholders
            has_title = False
            has_body = False
            for placeholder in layout.placeholders:
                if placeholder.placeholder_format.type == 1:  # Title
                    has_title = True
                elif placeholder.placeholder_format.type == 2:  # Body
                    has_body = True
            if has_title and has_body:
                slide_layout = layout
                break
        
        if slide_layout is None:
            # Fallback to layout index 1
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        
        # Add the slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Determine text color
        if text_color == "auto":
            # Auto-detect optimal color based on slide background
            r, g, b = get_optimal_text_color_for_slide(slide)
            color_rgb = RGBColor(r, g, b)
            print(f"  Auto-detected text color: #{r:02X}{g:02X}{b:02X}")
        else:
            # Parse specified color
            try:
                r = int(text_color[0:2], 16)
                g = int(text_color[2:4], 16)
                b = int(text_color[4:6], 16)
                color_rgb = RGBColor(r, g, b)
            except (ValueError, IndexError):
                color_rgb = RGBColor(0x33, 0x33, 0x33)  # Default dark gray
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            # Apply color
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = color_rgb
        
        # Find body placeholder and add items
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                tf = shape.text_frame
                tf.clear()  # Clear existing content
                
                for i, item in enumerate(items):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    
                    para.text = item
                    para.level = 0
                    
                    # Apply color
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                        run.font.size = Pt(18)
                
                break
        
        # Move slide to the specified position
        # Note: python-pptx adds slides at the end, we need to reorder
        if insert_position > 0 and insert_position < len(prs.slides):
            # Get the slide's XML element
            slide_id = prs.slides._sldIdLst[-1]  # Last added slide
            prs.slides._sldIdLst.remove(slide_id)
            prs.slides._sldIdLst.insert(insert_position, slide_id)
        
        print(f"  Added summary slide at position {insert_position + 1}")
        return True
        
    except Exception as e:
        print(f"  WARNING: Failed to add summary slide: {e}")
        return False


def main() -> None:
    """Main entry point for command-line usage."""
    args = parse_arguments()
    
    # Validate input file
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)
    
    # Validate replacements file
    replacements_path = Path(args.replacements)
    if not replacements_path.exists():
        print(f"Error: Replacements file not found: {args.replacements}")
        sys.exit(1)
    
    try:
        print(f"Loading presentation: {args.input}")
        prs = Presentation(input_path)
        
        print(f"Loading replacements: {args.replacements}")
        with open(replacements_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        # Extract slide management options
        slides_to_keep = data.pop("slides_to_keep", None)
        slides_to_delete = data.pop("slides_to_delete", None)
        add_summary = data.pop("add_summary_slide", False)
        summary_data = data.pop("summary_slide", None)
        
        # Remaining data is the replacements
        replacements = {}
        for key, value in data.items():
            if key.startswith("slide-"):
                replacements[key] = value
        
        print("Applying replacements...")
        auto_shrink = not args.no_auto_shrink
        modified, warnings = apply_replacements(prs, replacements, auto_shrink=auto_shrink)
        
        # Print warnings
        if warnings:
            print(f"\n=== Auto-adjustment Report ({len(warnings)} items) ===")
            for warning in warnings:
                print(warning)
            print("=" * 50 + "\n")
        
        # Delete slides if specified (AFTER applying replacements)
        deleted = 0
        if slides_to_keep is not None or slides_to_delete is not None:
            print("Processing slide deletion...")
            delete_data = {}
            if slides_to_keep is not None:
                delete_data["slides_to_keep"] = slides_to_keep
            if slides_to_delete is not None:
                delete_data["slides_to_delete"] = slides_to_delete
            deleted = delete_slides(prs, delete_data)
            print(f"  Deleted {deleted} slides")
        
        # Add summary slide if requested (AFTER slide deletion)
        if add_summary and summary_data:
            print("Adding summary slide...")
            add_summary_slide(prs, summary_data, insert_position=1)
        
        # Save output
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        
        total_shapes = sum(len(shapes) for shapes in replacements.values())
        print(f"Modified {modified} shapes across {len(replacements)} slides")
        if deleted > 0:
            print(f"Deleted {deleted} slides (final: {len(prs.slides)} slides)")
        print(f"Output saved to: {args.output}")
        
    except json.JSONDecodeError as e:
        print(f"Error parsing JSON: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error processing presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
