# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Classify input and determine optimal processing method.

This is a pure function script that replaces the Router agent.
No AI judgment needed - deterministic classification based on input characteristics.

Usage:
    python scripts/classify_input.py <input_path> [--output <classification.json>]
    python scripts/classify_input.py --detect-language <pptx_path>
    python scripts/classify_input.py --validate <classification.json>

Examples:
    python scripts/classify_input.py input/presentation.pptx
    python scripts/classify_input.py "https://qiita.com/example/items/xxx"
    python scripts/classify_input.py input/presentation.pptx --output output_manifest/classification.json
    python scripts/classify_input.py --validate output_manifest/classification.json

Output:
    JSON with input classification, recommended method, and base name.
    Follows workspace/classification.schema.json schema.

I/O Contract:
    Input: File path (PPTX/MD/JSON/TXT) or URL
    Output: classification.json following classification.schema.json
    Exit codes:
        0: Success
        1: Unknown input type or validation failed
        2: File not found
"""

import argparse
import json
import sys
import io
import re
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, Tuple

# Fix Windows console encoding issues
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# Language detection thresholds
ENGLISH_RATIO_THRESHOLD = 0.7  # >= 70% English characters = English PPTX

# Keyword mappings for base name generation
KEYWORD_MAPPINGS = {
    # Japanese to English
    "売上": "sales",
    "報告": "report",
    "進捗": "progress",
    "企画": "proposal",
    "提案": "proposal",
    "障害": "incident",
    "ブランチ": "branch",
    "ツール": "tool",
    "紹介": "intro",
    "新機能": "new_feature",
    "入門": "intro",
    "比較": "comparison",
    "やってみた": "trying",
    "まとめ": "summary",
    "解説": "explanation",
    "セキュリティ": "security",
    "データ": "data",
    "分析": "analysis",
    "設計": "design",
    "実装": "implementation",
}


def detect_language_ratio(text: str) -> float:
    """
    Calculate the ratio of English characters in the text.
    
    Returns:
        Float between 0 and 1 representing English character ratio.
    """
    if not text:
        return 0.0
    
    # Count ASCII letters vs total letters
    ascii_letters = sum(1 for c in text if c.isascii() and c.isalpha())
    total_letters = sum(1 for c in text if c.isalpha())
    
    if total_letters == 0:
        return 0.0
    
    return ascii_letters / total_letters


def detect_pptx_language(pptx_path: str) -> Tuple[str, float]:
    """
    Detect the primary language of a PPTX file.
    
    Returns:
        Tuple of (language_code, confidence)
        language_code: 'en' for English, 'ja' for Japanese, 'unknown'
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("Warning: python-pptx not installed. Cannot detect PPTX language.")
        return ("unknown", 0.0)
    
    prs = Presentation(pptx_path)
    all_text = []
    
    # Extract text from all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text)
    
    combined_text = " ".join(all_text)
    english_ratio = detect_language_ratio(combined_text)
    
    if english_ratio >= ENGLISH_RATIO_THRESHOLD:
        return ("en", english_ratio)
    elif english_ratio < 0.3:
        return ("ja", 1.0 - english_ratio)
    else:
        return ("mixed", 0.5)


def extract_keyword_from_title(title: str) -> str:
    """
    Extract English keyword from title for base name generation.
    """
    # Check for known Japanese mappings
    for jp_word, en_word in KEYWORD_MAPPINGS.items():
        if jp_word in title:
            return en_word
    
    # Try to extract English words
    english_words = re.findall(r'[A-Za-z][a-z]+', title)
    if english_words:
        # Use first 2-3 significant words
        keywords = [w.lower() for w in english_words if len(w) > 2][:3]
        if keywords:
            return "_".join(keywords)
    
    # Fallback to generic
    return "presentation"


def detect_input_type(input_path: str) -> Dict[str, Any]:
    """
    Detect the type of input and return classification.
    
    Args:
        input_path: File path or URL
        
    Returns:
        Dictionary with input_type, source, and metadata
    """
    # URL detection
    if input_path.startswith(('http://', 'https://')):
        # Detect specific platforms
        if 'qiita.com' in input_path:
            return {"input_type": "url", "platform": "qiita", "source": input_path}
        elif 'zenn.dev' in input_path:
            return {"input_type": "url", "platform": "zenn", "source": input_path}
        elif 'github.com' in input_path or 'github.io' in input_path:
            return {"input_type": "url", "platform": "github", "source": input_path}
        elif 'microsoft.com' in input_path or 'azure.com' in input_path:
            return {"input_type": "url", "platform": "microsoft", "source": input_path}
        else:
            return {"input_type": "url", "platform": "web", "source": input_path}
    
    # File detection
    path = Path(input_path)
    if not path.exists():
        return {"input_type": "unknown", "error": f"File not found: {input_path}"}
    
    suffix = path.suffix.lower()
    
    if suffix == '.pptx':
        lang, confidence = detect_pptx_language(input_path)
        return {
            "input_type": "pptx_en" if lang == "en" else "pptx_ja" if lang == "ja" else "pptx_mixed",
            "source_file": str(path),
            "language": lang,
            "language_confidence": confidence,
            "slide_count": get_slide_count(input_path)
        }
    elif suffix == '.md':
        return {"input_type": "markdown", "source_file": str(path)}
    elif suffix == '.json':
        return {"input_type": "json", "source_file": str(path)}
    elif suffix in ['.txt', '.text']:
        return {"input_type": "text", "source_file": str(path)}
    elif suffix in ['.png', '.jpg', '.jpeg', '.gif', '.webp']:
        return {"input_type": "image", "source_file": str(path)}
    else:
        return {"input_type": "unknown", "source_file": str(path), "extension": suffix}


def get_slide_count(pptx_path: str) -> int:
    """Get the number of slides in a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        return len(prs.slides)
    except Exception:
        return 0


def determine_method(input_info: Dict[str, Any]) -> Dict[str, Any]:
    """
    Determine the recommended processing method based on input classification.
    
    Returns:
        Dictionary with recommended_method, output_format, and reasoning
    """
    input_type = input_info.get("input_type", "unknown")
    
    if input_type == "pptx_en":
        return {
            "recommended_method": "reconstruct",
            "output_format": "template",
            "reasoning": "English PPTX detected - use reconstruct method to preserve slide master",
            "workflow": [
                "analyze_template.py",
                "extract_images.py", 
                "reconstruct_analyzer.py",
                "Localizer agent (translation)",
                "create_from_template.py"
            ]
        }
    
    elif input_type == "pptx_ja":
        return {
            "recommended_method": "analyze_only",
            "output_format": "template",
            "reasoning": "Japanese PPTX detected - may only need analysis or minor edits",
            "workflow": [
                "analyze_template.py",
                "reconstruct_analyzer.py"
            ]
        }
    
    elif input_type == "url":
        platform = input_info.get("platform", "web")
        return {
            "recommended_method": "create_new",
            "output_format": "template",
            "reasoning": f"Web content from {platform} - create new PPTX from content",
            "workflow": [
                "fetch_article (curl/API)",
                "extract_images (download)",
                "create content.json",
                "create_from_template.py"
            ]
        }
    
    elif input_type == "markdown":
        return {
            "recommended_method": "create_new",
            "output_format": "template",
            "reasoning": "Markdown input - convert to PPTX",
            "workflow": [
                "parse markdown",
                "create content.json",
                "create_from_template.py"
            ]
        }
    
    else:
        return {
            "recommended_method": "manual",
            "output_format": "unknown",
            "reasoning": f"Unknown input type: {input_type}",
            "workflow": []
        }


def generate_base_name(input_info: Dict[str, Any], purpose: str = "report") -> str:
    """
    Generate base name following the convention: {YYYYMMDD}_{keyword}_{purpose}
    """
    today = datetime.now().strftime("%Y%m%d")
    
    # Extract keyword from source
    source = input_info.get("source_file") or input_info.get("source", "")
    
    if source:
        path = Path(source.split('?')[0])  # Remove URL query params
        name = path.stem
        
        # Clean up the name
        keyword = extract_keyword_from_title(name)
    else:
        keyword = "presentation"
    
    # Sanitize keyword (max 30 chars, only a-z, 0-9, _)
    keyword = re.sub(r'[^a-z0-9_]', '_', keyword.lower())
    keyword = re.sub(r'_+', '_', keyword).strip('_')[:30]
    
    return f"{today}_{keyword}_{purpose}"


def classify_input(input_path: str, purpose: str = "report") -> Dict[str, Any]:
    """
    Main classification function.
    
    Returns complete context for Orchestrator.
    """
    # Step 1: Detect input type
    input_info = detect_input_type(input_path)
    
    # Step 2: Determine method
    method_info = determine_method(input_info)
    
    # Step 3: Generate base name
    base_name = generate_base_name(input_info, purpose)
    
    # Step 4: Build context
    context = {
        **input_info,
        **method_info,
        "base_name": base_name,
        "user_confirmation_required": True,
        "confidence": input_info.get("language_confidence", 0.8),
    }
    
    # Add confirmation prompt for PPTX
    if input_info.get("input_type", "").startswith("pptx"):
        slide_count = input_info.get("slide_count", 0)
        context["confirmation_prompt"] = f"""
英語 PPTX を検出しました ({slide_count} スライド)

処理方式を選択してください:
1. template（最推奨）: テンプレート継承で生成 ⭐⭐⭐⭐⭐
2. python-pptx: シンプルできれい ⭐⭐⭐⭐
3. pptxgenjs: コード向け ⭐⭐⭐⭐
"""
    
    return context


def validate_classification(classification_path: str) -> bool:
    """
    Validate a classification.json file against the schema.
    
    Returns True if valid, False otherwise.
    """
    try:
        import jsonschema
    except ImportError:
        print("Warning: jsonschema not installed. Skipping validation.")
        return True
    
    # Load schema
    schema_path = Path(__file__).parent.parent / "workspace" / "classification.schema.json"
    if not schema_path.exists():
        print(f"Warning: Schema not found at {schema_path}")
        return True
    
    with open(schema_path, "r", encoding="utf-8") as f:
        schema = json.load(f)
    
    # Load classification
    with open(classification_path, "r", encoding="utf-8") as f:
        classification = json.load(f)
    
    # Validate
    try:
        jsonschema.validate(instance=classification, schema=schema)
        print(f"✅ Classification validation passed: {classification_path}")
        return True
    except jsonschema.ValidationError as e:
        print(f"❌ Classification validation failed: {e.message}")
        print(f"   Path: {'.'.join(str(p) for p in e.absolute_path)}")
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Classify input and determine processing method"
    )
    parser.add_argument("input_path", nargs="?", help="Input file path or URL")
    parser.add_argument("--output", "-o", help="Output JSON file path")
    parser.add_argument("--purpose", "-p", default="report",
                        choices=["report", "lt", "incident", "blog", "custom"],
                        help="Purpose of the presentation")
    parser.add_argument("--detect-language", action="store_true",
                        help="Only detect language of PPTX")
    parser.add_argument("--validate", metavar="CLASSIFICATION_JSON",
                        help="Validate a classification.json file against schema")
    
    args = parser.parse_args()
    
    # Validate mode
    if args.validate:
        valid = validate_classification(args.validate)
        sys.exit(0 if valid else 1)
    
    # Require input_path for other modes
    if not args.input_path:
        parser.error("input_path is required")
    
    if args.detect_language:
        lang, conf = detect_pptx_language(args.input_path)
        print(f"Language: {lang} (confidence: {conf:.2%})")
        return
    
    # Classify input
    context = classify_input(args.input_path, args.purpose)
    
    # Output
    output_json = json.dumps(context, indent=2, ensure_ascii=False)
    
    if args.output:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(output_json)
        print(f"Classification saved to: {args.output}")
        
        # Validate output
        validate_classification(args.output)
    else:
        print(output_json)
    
    # Return appropriate exit code
    if context.get("input_type") == "unknown":
        sys.exit(1)


if __name__ == "__main__":
    main()
