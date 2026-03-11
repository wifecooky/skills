# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Clean PPTX template by removing problematic elements.

Usage:
    python scripts/clean_template.py <input.pptx> <output.pptx> [--keep-first-master]

Removes:
    - Background images (PICTURE shapes) from slide masters
    - Blip references from Picture Placeholders in layouts
    - Optionally keeps only the first slide master

This prepares external templates (e.g., from Microsoft Ignite) for use
without decorative background images interfering with generated content.
"""

from pptx import Presentation
from lxml import etree
import sys
from pathlib import Path
import shutil


def remove_master_background_images(prs) -> int:
    """Remove PICTURE shapes from slide masters."""
    removed = 0
    for master in prs.slide_masters:
        shapes_to_remove = []
        for shape in master.shapes:
            if 'PICTURE' in str(shape.shape_type):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
                removed += 1
            except Exception as e:
                print(f"  [!] Could not remove shape: {e}")
    
    return removed


def remove_placeholder_blips(prs) -> int:
    """Remove blip elements from Picture Placeholders in layouts."""
    removed = 0
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type)
                        if 'PICTURE' in ph_type:
                            # Find and remove blip elements
                            pic = shape._element
                            blip_elements = pic.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                            for blip in blip_elements:
                                blip.getparent().remove(blip)
                                removed += 1
                    except (ValueError, AttributeError):
                        pass
    
    return removed


def clean_template(input_path: str, output_path: str, keep_first_master_only: bool = False) -> dict:
    """
    Clean PPTX template and save to output path.
    
    Args:
        input_path: Path to input PPTX
        output_path: Path to output PPTX
        keep_first_master_only: If True, remove all but the first slide master
    
    Returns:
        dict with cleaning statistics
    """
    prs = Presentation(input_path)
    
    stats = {
        'input': Path(input_path).name,
        'output': Path(output_path).name,
        'masters_before': len(prs.slide_masters),
        'background_images_removed': 0,
        'placeholder_blips_removed': 0,
        'masters_removed': 0
    }
    
    print(f"\n🔧 Cleaning template: {input_path}")
    print(f"   Masters: {stats['masters_before']}")
    
    # Step 1: Remove background images from masters
    bg_removed = remove_master_background_images(prs)
    stats['background_images_removed'] = bg_removed
    print(f"   Removed {bg_removed} background image(s) from masters")
    
    # Step 2: Remove blip references from placeholders
    blip_removed = remove_placeholder_blips(prs)
    stats['placeholder_blips_removed'] = blip_removed
    print(f"   Removed {blip_removed} blip reference(s) from placeholders")
    
    # Step 3: Optionally remove extra masters (not fully implemented - complex)
    # This would require careful handling of slide relationships
    if keep_first_master_only and len(prs.slide_masters) > 1:
        print(f"   [i] Note: --keep-first-master not fully implemented yet")
    
    # Save cleaned template
    prs.save(output_path)
    
    total_removed = bg_removed + blip_removed
    stats['total_removed'] = total_removed
    
    print(f"\n✅ Saved clean template: {output_path}")
    print(f"   Total elements removed: {total_removed}")
    
    return stats


def main():
    if len(sys.argv) < 3:
        print("Usage: python scripts/clean_template.py <input.pptx> <output.pptx> [--keep-first-master]")
        print()
        print("Options:")
        print("  --keep-first-master  Remove all but the first slide master (experimental)")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    keep_first_master = '--keep-first-master' in sys.argv
    
    if not Path(input_path).exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    # Create output directory if needed
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    
    stats = clean_template(input_path, output_path, keep_first_master)
    
    # Exit code: 0 on success
    sys.exit(0)


if __name__ == "__main__":
    main()
