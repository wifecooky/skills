# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
#
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Create a clean template from an existing PPTX file.

This script extracts the slide master and layouts from a source PPTX,
removes problematic elements, and creates a clean template for reuse.

Usage:
    python scripts/create_clean_template.py <source.pptx> <output_template.pptx> [options]

Options:
    --remove-backgrounds    Remove background images from layouts
    --fix-placeholders      Adjust placeholder positions for consistency
    --remove-decorations    Remove decorative shapes (rectangles, lines, etc.)
    --analyze               Analyze template without creating (dry run)
    --all                   Apply all fixes (default)

Example:
    python scripts/create_clean_template.py input/BRK252.pptx templates/brk252_clean.pptx --all
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Optional, List, Dict, Any
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.dml.color import RGBColor
    from lxml import etree
except ImportError as e:
    print(f"Error: Required package not found: {e}")
    print("Install with: pip install python-pptx lxml")
    sys.exit(1)


def analyze_layout(layout, slide_height: float) -> Dict[str, Any]:
    """Analyze a single layout for potential issues."""
    issues = []
    info = {
        "name": layout.name,
        "placeholders": [],
        "background_images": [],
        "decorative_shapes": [],
        "issues": issues
    }
    
    # Check placeholders
    for ph in layout.placeholders:
        ph_info = {
            "idx": ph.placeholder_format.idx,
            "type": str(ph.placeholder_format.type),
            "left": ph.left.inches,
            "top": ph.top.inches,
            "width": ph.width.inches,
            "height": ph.height.inches,
            "top_percent": ph.top.inches / slide_height * 100
        }
        info["placeholders"].append(ph_info)
        
        # Check for problematic positioning
        if ph_info["top_percent"] > 70:
            issues.append(f"Placeholder {ph_info['type']} is positioned very low ({ph_info['top_percent']:.0f}%)")
        if ph_info["width"] < 5 and "TITLE" in ph_info["type"]:
            issues.append(f"Title placeholder is narrow ({ph_info['width']:.1f} inches)")
    
    # Check non-placeholder shapes
    for shape in layout.shapes:
        if shape.is_placeholder:
            continue
        
        shape_info = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left.inches,
            "width": shape.width.inches
        }
        
        # Background images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info["background_images"].append(shape_info)
            issues.append(f"Background image found: {shape.name}")
        
        # Decorative shapes (rectangles, lines at edges)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.left.inches < 0.5 or shape.width.inches > 10:
                info["decorative_shapes"].append(shape_info)
                issues.append(f"Decorative shape at edge: {shape.name}")
    
    return info


def analyze_template(prs: Presentation) -> Dict[str, Any]:
    """Analyze entire template for issues."""
    slide_height = prs.slide_height.inches
    slide_width = prs.slide_width.inches
    
    analysis = {
        "slide_size": f"{slide_width:.2f} x {slide_height:.2f} inches",
        "layout_count": len(prs.slide_layouts),
        "layouts": [],
        "summary": {
            "total_issues": 0,
            "layouts_with_background_images": 0,
            "layouts_with_decorations": 0,
            "layouts_with_position_issues": 0
        }
    }
    
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = analyze_layout(layout, slide_height)
        layout_info["index"] = i
        analysis["layouts"].append(layout_info)
        
        # Update summary
        if layout_info["background_images"]:
            analysis["summary"]["layouts_with_background_images"] += 1
        if layout_info["decorative_shapes"]:
            analysis["summary"]["layouts_with_decorations"] += 1
        if any("position" in issue.lower() for issue in layout_info["issues"]):
            analysis["summary"]["layouts_with_position_issues"] += 1
        analysis["summary"]["total_issues"] += len(layout_info["issues"])
    
    return analysis


def remove_background_images(prs: Presentation) -> int:
    """Remove background images from all layouts."""
    removed_count = 0
    
    for layout in prs.slide_layouts:
        shapes_to_remove = []
        for shape in layout.shapes:
            if not shape.is_placeholder and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            removed_count += 1
    
    # Also check slide masters
    for master in prs.slide_masters:
        shapes_to_remove = []
        for shape in master.shapes:
            if not shape.is_placeholder and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            removed_count += 1
    
    return removed_count


def remove_decorative_shapes(prs: Presentation) -> int:
    """Remove decorative shapes (edge rectangles, lines, out-of-bounds textboxes, etc.)."""
    removed_count = 0
    slide_width = prs.slide_width.inches
    
    for layout in prs.slide_layouts:
        shapes_to_remove = []
        for shape in layout.shapes:
            if shape.is_placeholder:
                continue
            
            # Remove shapes outside slide bounds (causes vertical text issues)
            if shape.left.inches > slide_width:
                shapes_to_remove.append(shape)
                continue
            
            # Remove textboxes at edges that might cause issues
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # TextBox outside visible area
                if shape.left.inches > slide_width * 0.95:
                    shapes_to_remove.append(shape)
                    continue
                # Very narrow textbox at left edge (vertical text issue)
                if shape.left.inches < 0.3 and shape.width.inches < 0.5:
                    shapes_to_remove.append(shape)
                    continue
            
            # Remove decorative auto-shapes at edges
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Left edge decorations (vertical bars, etc.)
                if shape.left.inches < 0.3 and shape.width.inches < 1.0:
                    shapes_to_remove.append(shape)
                # Full-width decorations
                elif shape.width.inches > slide_width * 0.9:
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            removed_count += 1
    
    return removed_count


def has_dark_background(layout) -> bool:
    """Check if a layout has a dark background (image or dark fill)."""
    for shape in layout.shapes:
        if not shape.is_placeholder:
            # Check for background images (usually dark in presentation templates)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
    
    # Check background fill
    try:
        bg = layout.background
        if bg.fill.type is not None:
            # Check if it's a solid dark fill
            if hasattr(bg.fill, 'fore_color') and bg.fill.fore_color.rgb:
                rgb = bg.fill.fore_color.rgb
                # Calculate luminance (simple approximation)
                luminance = (int(str(rgb)[0:2], 16) * 0.299 + 
                           int(str(rgb)[2:4], 16) * 0.587 + 
                           int(str(rgb)[4:6], 16) * 0.114)
                if luminance < 128:
                    return True
    except:
        pass
    
    return False


def fix_text_contrast(prs: Presentation) -> int:
    """Fix text colors for better contrast on dark/light backgrounds."""
    fixed_count = 0
    
    for layout in prs.slide_layouts:
        is_dark = has_dark_background(layout)
        
        if is_dark:
            # Dark background: ensure text is light colored
            for ph in layout.placeholders:
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        # Check if font color is dark (needs fixing)
                        try:
                            if para.font.color.rgb:
                                rgb = para.font.color.rgb
                                luminance = (int(str(rgb)[0:2], 16) * 0.299 + 
                                           int(str(rgb)[2:4], 16) * 0.587 + 
                                           int(str(rgb)[4:6], 16) * 0.114)
                                if luminance < 180:  # Dark text on dark background
                                    para.font.color.rgb = RGBColor(255, 255, 255)
                                    fixed_count += 1
                        except:
                            pass
    
    return fixed_count


def fix_placeholder_positions(prs: Presentation) -> int:
    """Adjust placeholder positions for better layout.
    
    Conservative approach: only fix clearly problematic positions,
    don't change layouts that might be intentionally designed.
    """
    fixed_count = 0
    slide_height = prs.slide_height.inches
    slide_width = prs.slide_width.inches
    
    for layout in prs.slide_layouts:
        layout_name = layout.name.lower()
        is_section_layout = 'section' in layout_name
        is_title_layout = layout_name.startswith('title') and 'only' not in layout_name
        
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            
            # Fix title placeholders - ONLY for section/title layouts
            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                # Only adjust section slide titles that are positioned oddly
                if is_section_layout:
                    top_percent = ph.top.inches / slide_height
                    # Only fix if really out of range (< 10% or > 60%)
                    if top_percent < 0.10 or top_percent > 0.60:
                        # IMPORTANT: Preserve width/height when adjusting top. See TROUBLESHOOTING.md #48.
                        original_width = ph.width
                        original_height = ph.height
                        ph.top = Inches(slide_height * 0.35)
                        ph.width = original_width
                        ph.height = original_height
                        fixed_count += 1
                
                # For title slides, ensure title is not too narrow
                if is_title_layout:
                    if ph.width.inches < slide_width * 0.5:
                        ph.width = Inches(slide_width * 0.75)
                        # Center the title
                        ph.left = Inches((slide_width - ph.width.inches) / 2)
                        fixed_count += 1
            
            # Fix subtitle placeholders - only if really far from title
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                if ph.top.inches > slide_height * 0.7:
                    # IMPORTANT: Preserve width/height when adjusting top. See TROUBLESHOOTING.md #48.
                    original_width = ph.width
                    original_height = ph.height
                    ph.top = Inches(slide_height * 0.55)
                    ph.width = original_width
                    ph.height = original_height
                    fixed_count += 1
            
            # Fix body/content placeholders - only negative positions
            elif ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                if ph.left.inches < 0:
                    ph.left = Inches(0.5)
                    fixed_count += 1
    
    return fixed_count


def create_clean_template(
    source_path: str,
    output_path: str,
    remove_backgrounds: bool = False,  # ★ Keep backgrounds by default
    remove_decorations: bool = True,
    fix_positions: bool = True,
    fix_contrast: bool = True  # ★ NEW: Fix text contrast
) -> Dict[str, Any]:
    """Create a clean template from source PPTX."""
    
    print(f"Loading source: {source_path}")
    prs = Presentation(source_path)
    
    results = {
        "source": source_path,
        "output": output_path,
        "slide_size": f"{prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}",
        "original_layouts": len(prs.slide_layouts),
        "actions": []
    }
    
    # Remove existing slides (keep only layouts)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    results["actions"].append("Removed all existing slides")
    
    # Remove background images
    if remove_backgrounds:
        count = remove_background_images(prs)
        if count > 0:
            results["actions"].append(f"Removed {count} background images")
    
    # Remove decorative shapes
    if remove_decorations:
        count = remove_decorative_shapes(prs)
        if count > 0:
            results["actions"].append(f"Removed {count} decorative shapes")
    
    # Fix placeholder positions
    if fix_positions:
        count = fix_placeholder_positions(prs)
        if count > 0:
            results["actions"].append(f"Fixed {count} placeholder positions")
    
    # Fix text contrast for dark backgrounds
    if fix_contrast:
        count = fix_text_contrast(prs)
        if count > 0:
            results["actions"].append(f"Fixed {count} text colors for better contrast")
    
    # Save the clean template
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    
    results["success"] = True
    print(f"\n✅ Created clean template: {output_path}")
    
    return results


def main():
    parser = argparse.ArgumentParser(
        description='Create a clean template from an existing PPTX file.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Analyze template without changes
  python scripts/create_clean_template.py input/BRK252.pptx --analyze

  # Create clean template with all fixes
  python scripts/create_clean_template.py input/BRK252.pptx templates/brk252_clean.pptx --all

  # Create template with specific fixes
  python scripts/create_clean_template.py input/BRK252.pptx templates/out.pptx --remove-backgrounds
        """
    )
    
    parser.add_argument('source', help='Source PPTX file')
    parser.add_argument('output', nargs='?', help='Output template path')
    parser.add_argument('--analyze', action='store_true', help='Analyze only (dry run)')
    parser.add_argument('--remove-backgrounds', action='store_true', help='Remove background images')
    parser.add_argument('--remove-decorations', action='store_true', help='Remove decorative shapes')
    parser.add_argument('--fix-placeholders', action='store_true', help='Fix placeholder positions')
    parser.add_argument('--all', action='store_true', help='Apply all fixes (default)')
    parser.add_argument('--json', action='store_true', help='Output analysis as JSON')
    
    args = parser.parse_args()
    
    # Check source exists
    if not Path(args.source).exists():
        print(f"Error: Source file not found: {args.source}")
        sys.exit(1)
    
    # Analyze mode
    if args.analyze:
        prs = Presentation(args.source)
        analysis = analyze_template(prs)
        
        if args.json:
            print(json.dumps(analysis, indent=2, ensure_ascii=False))
        else:
            print(f"\n{'='*60}")
            print(f"Template Analysis: {args.source}")
            print(f"{'='*60}")
            print(f"Slide size: {analysis['slide_size']}")
            print(f"Total layouts: {analysis['layout_count']}")
            print(f"\nSummary:")
            print(f"  Total issues: {analysis['summary']['total_issues']}")
            print(f"  Layouts with background images: {analysis['summary']['layouts_with_background_images']}")
            print(f"  Layouts with decorations: {analysis['summary']['layouts_with_decorations']}")
            print(f"  Layouts with position issues: {analysis['summary']['layouts_with_position_issues']}")
            
            # Show problematic layouts
            print(f"\nProblematic layouts:")
            for layout in analysis['layouts']:
                if layout['issues']:
                    print(f"\n  [{layout['index']}] {layout['name']}:")
                    for issue in layout['issues']:
                        print(f"      ⚠️ {issue}")
        
        sys.exit(0)
    
    # Create mode - require output path
    if not args.output:
        print("Error: Output path required (or use --analyze)")
        sys.exit(1)
    
    # Determine which fixes to apply
    if args.all:
        # --all: apply decorations removal and position fixes (NOT backgrounds)
        remove_backgrounds = False  # ★ Keep backgrounds by default
        remove_decorations = True
        fix_positions = True
    elif not (args.remove_backgrounds or args.remove_decorations or args.fix_placeholders):
        # No specific options: apply safe defaults
        remove_backgrounds = False
        remove_decorations = True
        fix_positions = True
    else:
        remove_backgrounds = args.remove_backgrounds
        remove_decorations = args.remove_decorations
        fix_positions = args.fix_placeholders
    
    # Create clean template
    results = create_clean_template(
        args.source,
        args.output,
        remove_backgrounds=remove_backgrounds,
        remove_decorations=remove_decorations,
        fix_positions=fix_positions
    )
    
    # Print results
    print(f"\nActions performed:")
    for action in results['actions']:
        print(f"  ✓ {action}")
    
    print(f"\nNext steps:")
    print(f"  1. Analyze layouts: python scripts/analyze_template.py {args.output}")
    print(f"  2. Generate PPTX: python scripts/create_from_template.py {args.output} content.json output.pptx")


if __name__ == "__main__":
    main()
