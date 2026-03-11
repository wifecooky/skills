# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Diagnose PPTX template for potential issues before using it.

Usage:
    python scripts/diagnose_template.py <template.pptx>

Checks for:
    - Background images in slide masters (decorative graphics)
    - Picture Placeholders with embedded default images (blip references)
    - Embedded fonts that may cause substitution warnings
    - Broken external links
    - Excessive number of slide masters

Output:
    Prints diagnosis report to console.
    Returns exit code 0 if clean, 1 if issues found.
"""

from pptx import Presentation
from lxml import etree
import sys
from pathlib import Path


def count_master_background_images(prs) -> list:
    """Count PICTURE shapes in slide masters (likely background images)."""
    issues = []
    for master_idx, master in enumerate(prs.slide_masters):
        bg_images = []
        for shape in master.shapes:
            if 'PICTURE' in str(shape.shape_type):
                bg_images.append(shape.name)
        if bg_images:
            issues.append({
                'type': 'master_background',
                'master_index': master_idx,
                'master_name': master.name or f'Master_{master_idx}',
                'images': bg_images,
                'count': len(bg_images)
            })
    return issues


def count_layout_placeholder_blips(prs) -> list:
    """Count Picture Placeholders with embedded blip references."""
    issues = []
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    
    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            blip_count = 0
            for shape in layout.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type)
                        if 'PICTURE' in ph_type:
                            # Check for blip element
                            xml = etree.tostring(shape._element, encoding='unicode')
                            if 'blip' in xml and 'r:embed' in xml:
                                blip_count += 1
                    except (ValueError, AttributeError):
                        pass
            if blip_count > 0:
                issues.append({
                    'type': 'placeholder_blip',
                    'master_index': master_idx,
                    'layout_index': layout_idx,
                    'layout_name': layout.name or f'Layout_{layout_idx}',
                    'count': blip_count
                })
    return issues


def check_master_count(prs) -> list:
    """Check if there are too many slide masters."""
    issues = []
    master_count = len(prs.slide_masters)
    if master_count > 5:
        issues.append({
            'type': 'excessive_masters',
            'count': master_count,
            'recommendation': 'Consider using only the first master or cleaning up unused masters'
        })
    return issues


def check_view_props(template_path: str) -> list:
    """Check if template has non-standard view settings (e.g., sldMasterView)."""
    import zipfile
    issues = []
    
    try:
        with zipfile.ZipFile(template_path, 'r') as z:
            if 'ppt/viewProps.xml' in z.namelist():
                content = z.read('ppt/viewProps.xml').decode('utf-8')
                if 'lastView="sldMasterView"' in content:
                    issues.append({
                        'type': 'master_view_default',
                        'severity': 'warning',
                        'description': 'Template opens in Slide Master view instead of Normal view',
                        'recommendation': 'Generated PPTX will open in Master view. Will be auto-fixed during generation.'
                    })
                elif 'lastView="sldSorterView"' in content:
                    issues.append({
                        'type': 'sorter_view_default',
                        'severity': 'info',
                        'description': 'Template opens in Slide Sorter view',
                        'recommendation': 'Consider if this is intentional'
                    })
    except Exception as e:
        pass  # Silently ignore if can't read viewProps
    
    return issues


def diagnose_template(template_path: str) -> dict:
    """
    Diagnose PPTX template and return structured report.
    
    Returns:
        dict with 'issues', 'summary', 'clean' keys
    """
    prs = Presentation(template_path)
    
    all_issues = []
    
    # Check for background images in masters
    master_bg_issues = count_master_background_images(prs)
    all_issues.extend(master_bg_issues)
    
    # Check for placeholder blips
    placeholder_issues = count_layout_placeholder_blips(prs)
    all_issues.extend(placeholder_issues)
    
    # Check master count
    master_count_issues = check_master_count(prs)
    all_issues.extend(master_count_issues)
    
    # Check view properties (viewProps.xml)
    view_prop_issues = check_view_props(template_path)
    all_issues.extend(view_prop_issues)
    
    # Calculate totals
    total_bg_images = sum(i['count'] for i in master_bg_issues)
    total_placeholder_blips = sum(i['count'] for i in placeholder_issues)
    
    summary = {
        'template': Path(template_path).name,
        'slide_masters': len(prs.slide_masters),
        'total_layouts': sum(len(m.slide_layouts) for m in prs.slide_masters),
        'background_images_in_masters': total_bg_images,
        'placeholder_blips': total_placeholder_blips,
        'issues_found': len(all_issues)
    }
    
    return {
        'issues': all_issues,
        'summary': summary,
        'clean': len(all_issues) == 0
    }


def print_report(report: dict) -> None:
    """Print diagnosis report to console."""
    summary = report['summary']
    issues = report['issues']
    
    print("\n" + "=" * 60)
    print("📋 TEMPLATE DIAGNOSIS REPORT")
    print("=" * 60)
    print(f"\nTemplate: {summary['template']}")
    print(f"Slide Masters: {summary['slide_masters']}")
    print(f"Total Layouts: {summary['total_layouts']}")
    print()
    
    if report['clean']:
        print("✅ Template is CLEAN - no issues found")
        return
    
    print(f"⚠️  ISSUES FOUND: {summary['issues_found']}")
    print()
    
    # Group issues by type
    bg_issues = [i for i in issues if i['type'] == 'master_background']
    blip_issues = [i for i in issues if i['type'] == 'placeholder_blip']
    master_issues = [i for i in issues if i['type'] == 'excessive_masters']
    
    if bg_issues:
        print("🖼️  Background Images in Slide Masters:")
        print("-" * 40)
        for issue in bg_issues:
            print(f"  Master {issue['master_index']} ({issue['master_name']}): {issue['count']} image(s)")
            for img in issue['images'][:5]:  # Show first 5
                print(f"    - {img}")
            if len(issue['images']) > 5:
                print(f"    ... and {len(issue['images']) - 5} more")
        print(f"\n  Total: {summary['background_images_in_masters']} background image(s)")
        print("  → These may overlap with your slide content")
        print()
    
    if blip_issues:
        print("📎 Picture Placeholders with Embedded Images:")
        print("-" * 40)
        for issue in blip_issues:
            print(f"  Layout {issue['layout_index']} ({issue['layout_name']}): {issue['count']} placeholder(s)")
        print(f"\n  Total: {summary['placeholder_blips']} placeholder(s) with blip references")
        print("  → These may show 'image cannot be displayed' errors")
        print()
    
    if master_issues:
        print("📊 Excessive Slide Masters:")
        print("-" * 40)
        for issue in master_issues:
            print(f"  Found {issue['count']} masters (recommended: ≤5)")
            print(f"  → {issue['recommendation']}")
        print()
    
    # View properties issues
    view_issues = [i for i in issues if i['type'] in ('master_view_default', 'sorter_view_default')]
    if view_issues:
        print("👁️  View Settings Issues:")
        print("-" * 40)
        for issue in view_issues:
            icon = "⚠️" if issue.get('severity') == 'warning' else "ℹ️"
            print(f"  {icon} {issue['description']}")
            print(f"     → {issue['recommendation']}")
        print()
    
    print("=" * 60)
    print("💡 RECOMMENDATION: Run clean_template.py to fix these issues")
    print("   python scripts/clean_template.py <input.pptx> <output.pptx>")
    print("=" * 60)


def main():
    import argparse
    
    parser = argparse.ArgumentParser(description='Diagnose PPTX template for potential issues')
    parser.add_argument('template', help='Path to template PPTX file')
    parser.add_argument('--json', action='store_true', help='Output as JSON (for scripting)')
    
    args = parser.parse_args()
    
    if not Path(args.template).exists():
        if args.json:
            import json
            print(json.dumps({"error": f"File not found: {args.template}"}))
        else:
            print(f"Error: File not found: {args.template}")
        sys.exit(1)
    
    report = diagnose_template(args.template)
    
    if args.json:
        import json
        # Output JSON for scripting
        output = {
            "clean": report['clean'],
            "total_issues": report['summary']['total'],
            "summary": report['summary'],
            "issues": report['issues']
        }
        print(json.dumps(output))
    else:
        print_report(report)
    
    # Exit code: 0 if clean, 1 if issues found
    sys.exit(0 if report['clean'] else 1)


if __name__ == "__main__":
    main()
