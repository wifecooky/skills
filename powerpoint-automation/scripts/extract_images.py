# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Extract images from PPTX slides.

Usage:
    python scripts/extract_images.py <input.pptx> <output_dir> [--skip-icons]

Creates images/slide_01.png, images/slide_02.png, etc. for each slide.

Options:
    --skip-icons    Skip small images that are likely icons/logos (< 400px or square < 800px)
"""

import sys
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def is_icon_or_logo(width_px: int, height_px: int, min_size: int = 400) -> bool:
    """
    Detect if image is likely an icon/logo.
    
    Criteria:
    - Very small size (< min_size on any dimension)
    - Square aspect ratio (0.9-1.1) AND small size (<= 800px)
    
    Note: Icons/logos are now extracted but marked, not skipped.
    """
    if width_px < min_size or height_px < min_size:
        return True
    
    aspect_ratio = width_px / height_px if height_px > 0 else 1
    if 0.9 <= aspect_ratio <= 1.1 and max(width_px, height_px) <= 800:
        return True
    
    return False


def extract_images_from_pptx(pptx_path: str, output_dir: str, skip_icons: bool = False, mark_icons: bool = True) -> None:
    """
    Extract the main image from each slide.
    
    Args:
        pptx_path: Path to the input PPTX file.
        output_dir: Directory to save extracted images.
        skip_icons: If True, skip icon/logo images entirely.
        mark_icons: If True, log when icons/logos are detected (but still extract them).
    """
    # Prevent output to input/ directory
    out_path = Path(output_dir)
    if out_path.parts and out_path.parts[0] == 'input':
        print("Error: Cannot output to input/ directory. Use images/ or output_manifest/ instead.")
        sys.exit(1)
    
    prs = Presentation(pptx_path)
    out_path.mkdir(parents=True, exist_ok=True)
    
    print(f"\nExtracting images from: {pptx_path}")
    print(f"Output directory: {output_dir}")
    if skip_icons:
        print("Skipping icons/logos: YES")
    print("=" * 50)
    
    # Also extract from the zip structure for embedded images
    images_extracted = 0
    icons_skipped = 0
    
    with zipfile.ZipFile(pptx_path, 'r') as z:
        media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
        print(f"Found {len(media_files)} embedded media files")
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        
        # Find images on this slide
        images_on_slide = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Get image data
                    image = shape.image
                    ext = image.ext
                    blob = image.blob
                    
                    images_on_slide.append({
                        'ext': ext,
                        'blob': blob,
                        'width': shape.width,
                        'height': shape.height,
                        'left': shape.left,
                        'top': shape.top
                    })
                except Exception as e:
                    print(f"  [!] Error extracting image from slide {slide_num}: {e}")
        
        if images_on_slide:
            # Filter out icons/logos if requested
            if skip_icons:
                content_images = []
                for img in images_on_slide:
                    # Convert EMU to pixels (1 inch = 914400 EMU, 96 DPI)
                    width_px = int(img['width'] / 914400 * 96)
                    height_px = int(img['height'] / 914400 * 96)
                    
                    if is_icon_or_logo(width_px, height_px):
                        icons_skipped += 1
                        print(f"  Slide {slide_num:3d}: Skipped icon/logo ({width_px}x{height_px}px)")
                    else:
                        content_images.append(img)
                
                images_on_slide = content_images
            
            if not images_on_slide:
                continue
                
            # Save the largest image (main content image)
            # Prefer landscape/screenshot images (16:9 or similar)
            def image_score(img):
                """Score images by size and aspect ratio (prefer 16:9 screenshots)."""
                width_px = int(img['width'] / 914400 * 96)
                height_px = int(img['height'] / 914400 * 96)
                area = width_px * height_px
                
                # Prefer 16:9 aspect ratio (1.7-1.8)
                aspect = width_px / height_px if height_px > 0 else 1
                aspect_bonus = 1.5 if 1.6 <= aspect <= 1.9 else 1.0
                
                return area * aspect_bonus
            
            main_image = max(images_on_slide, key=image_score)
            
            output_file = out_path / f"slide_{slide_num:02d}.{main_image['ext']}"
            with open(output_file, 'wb') as f:
                f.write(main_image['blob'])
            
            width_px = int(main_image['width'] / 914400 * 96)
            height_px = int(main_image['height'] / 914400 * 96)
            
            # Mark if image is icon/logo (for reference)
            is_icon = is_icon_or_logo(width_px, height_px)
            icon_marker = " 🔹icon/logo" if is_icon and mark_icons else ""
            
            print(f"  Slide {slide_num:3d}: Saved ({width_px}x{height_px}px) → {output_file.name}{icon_marker}")
            images_extracted += 1
    
    print("=" * 50)
    print(f"✅ Extracted images from {images_extracted} slides to {output_dir}")
    if skip_icons:
        print(f"   Skipped {icons_skipped} icon/logo images")


def main():
    if len(sys.argv) < 3:
        print("Usage: python extract_images.py <input.pptx> <output_dir> [--skip-icons]")
        sys.exit(1)
    
    skip_icons = '--skip-icons' in sys.argv
    extract_images_from_pptx(sys.argv[1], sys.argv[2], skip_icons=skip_icons)


if __name__ == "__main__":
    main()
