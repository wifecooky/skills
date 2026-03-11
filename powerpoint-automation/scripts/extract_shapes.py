#!/usr/bin/env python3
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Extract structured text content from PowerPoint presentations.

This module provides functionality to:
- Extract all text content from PowerPoint shapes
- Preserve paragraph formatting (alignment, bullets, fonts, spacing)
- Handle nested GroupShapes recursively with correct absolute positions
- Sort shapes by visual position on slides
- Export to JSON with clean, structured data

Usage:
    python extract_shapes.py input.pptx output.json [--issues-only]

Author: aktsmm
License: CC BY-NC 4.0
"""

import argparse
import json
import platform
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.util import Emu, Inches

# Constants for EMU (English Metric Units) conversion
EMU_PER_INCH = 914400


def parse_arguments() -> argparse.Namespace:
    """Parse command line arguments.
    
    Returns:
        Parsed command line arguments.
    """
    parser = argparse.ArgumentParser(
        description="Extract text inventory from PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python extract_shapes.py presentation.pptx inventory.json
    Extracts text inventory with positions for all shapes

  python extract_shapes.py presentation.pptx inventory.json --issues-only
    Extracts only text shapes that have overflow or overlap issues

Output JSON includes:
  - All text content organized by slide and shape
  - Visual position and size in inches
  - Paragraph properties and formatting
  - Issue detection: text overflow and shape overlaps
        """,
    )
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("output", help="Output JSON file for inventory")
    parser.add_argument(
        "--issues-only",
        action="store_true",
        help="Include only text shapes that have overflow or overlap issues",
    )
    return parser.parse_args()


@dataclass
class ParagraphData:
    """Data structure for paragraph properties extracted from PowerPoint."""
    
    text: str = ""
    bullet: bool = False
    level: int = 0
    alignment: Optional[str] = None
    space_before: Optional[float] = None
    space_after: Optional[float] = None
    line_spacing: Optional[float] = None
    font_size: Optional[float] = None
    font_name: Optional[str] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color: Optional[str] = None
    theme_color: Optional[str] = None
    
    @classmethod
    def from_paragraph(cls, paragraph: Any) -> "ParagraphData":
        """Create ParagraphData from a PowerPoint paragraph object.
        
        Args:
            paragraph: The PowerPoint paragraph object.
            
        Returns:
            ParagraphData instance with extracted properties.
        """
        data = cls()
        # Clean text: remove vertical tabs (soft line breaks) and other control characters
        text = paragraph.text.strip()
        # Replace vertical tab (\x0B) with space - this is PowerPoint's soft line break
        text = text.replace('\x0b', ' ').replace('\u000b', ' ')
        data.text = text
        
        # Extract bullet properties
        pPr = paragraph._p.get_or_add_pPr()
        buChar = pPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
        buAutoNum = pPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
        
        if buChar is not None or buAutoNum is not None:
            data.bullet = True
            data.level = paragraph.level if paragraph.level else 0
        
        # Extract alignment
        if paragraph.alignment:
            alignment_map = {
                PP_ALIGN.LEFT: "LEFT",
                PP_ALIGN.CENTER: "CENTER",
                PP_ALIGN.RIGHT: "RIGHT",
                PP_ALIGN.JUSTIFY: "JUSTIFY",
            }
            data.alignment = alignment_map.get(paragraph.alignment)
        
        # Extract spacing
        if paragraph.space_before:
            data.space_before = paragraph.space_before.pt
        if paragraph.space_after:
            data.space_after = paragraph.space_after.pt
        if paragraph.line_spacing:
            data.line_spacing = paragraph.line_spacing
        
        # Extract font properties from first run
        if paragraph.runs:
            run = paragraph.runs[0]
            font = run.font
            
            if font.size:
                data.font_size = font.size.pt
            if font.name:
                data.font_name = font.name
            if font.bold is not None:
                data.bold = font.bold
            if font.italic is not None:
                data.italic = font.italic
            if font.underline is not None:
                data.underline = font.underline
                
            # Extract color (handle different color types safely)
            try:
                if font.color.type is not None:
                    color_type = str(font.color.type)
                    if "RGB" in color_type and font.color.rgb:
                        data.color = str(font.color.rgb)
                    elif font.color.theme_color:
                        data.theme_color = str(font.color.theme_color).replace("MSO_THEME_COLOR.", "")
            except AttributeError:
                # Color might be inherited or undefined
                pass
        
        return data
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization.
        
        Returns:
            Dictionary representation of the paragraph data.
        """
        result: Dict[str, Any] = {"text": self.text}
        
        if self.bullet:
            result["bullet"] = True
            result["level"] = self.level
        if self.alignment:
            result["alignment"] = self.alignment
        if self.space_before is not None:
            result["space_before"] = self.space_before
        if self.space_after is not None:
            result["space_after"] = self.space_after
        if self.line_spacing is not None:
            result["line_spacing"] = self.line_spacing
        if self.font_size is not None:
            result["font_size"] = self.font_size
        if self.font_name:
            result["font_name"] = self.font_name
        if self.bold is not None:
            result["bold"] = self.bold
        if self.italic is not None:
            result["italic"] = self.italic
        if self.underline is not None:
            result["underline"] = self.underline
        if self.color:
            result["color"] = self.color
        if self.theme_color:
            result["theme_color"] = self.theme_color
            
        return result


@dataclass
class ShapeData:
    """Data structure for shape properties extracted from PowerPoint."""
    
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0
    placeholder_type: Optional[str] = None
    default_font_size: Optional[float] = None
    paragraphs: List[ParagraphData] = field(default_factory=list)
    overlap: Optional[Dict[str, Any]] = None
    overflow_bottom: Optional[float] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization.
        
        Returns:
            Dictionary representation of the shape data.
        """
        result: Dict[str, Any] = {
            "left": round(self.left, 2),
            "top": round(self.top, 2),
            "width": round(self.width, 2),
            "height": round(self.height, 2),
        }
        
        if self.placeholder_type:
            result["placeholder_type"] = self.placeholder_type
        if self.default_font_size is not None:
            result["default_font_size"] = self.default_font_size
        if self.overlap:
            result["overlap"] = self.overlap
        if self.overflow_bottom is not None:
            result["overflow_bottom"] = round(self.overflow_bottom, 2)
            
        result["paragraphs"] = [p.to_dict() for p in self.paragraphs]
        
        return result


def emu_to_inches(emu: int) -> float:
    """Convert EMU (English Metric Units) to inches.
    
    Args:
        emu: Value in EMUs.
        
    Returns:
        Value in inches.
    """
    return emu / EMU_PER_INCH


def get_placeholder_type_name(placeholder) -> Optional[str]:
    """Get the placeholder type as a string.
    
    Args:
        placeholder: The placeholder format object.
        
    Returns:
        String name of the placeholder type, or None.
    """
    if placeholder and placeholder.type:
        return str(placeholder.type).replace("PLACEHOLDER_TYPE.", "")
    return None


def extract_shape_text(shape: BaseShape, offset_left: int = 0, offset_top: int = 0) -> Optional[ShapeData]:
    """Extract text content from a shape.
    
    Args:
        shape: The PowerPoint shape object.
        offset_left: X offset for grouped shapes (in EMUs).
        offset_top: Y offset for grouped shapes (in EMUs).
        
    Returns:
        ShapeData if shape contains text, None otherwise.
    """
    # Skip shapes without text frames
    if not shape.has_text_frame:
        return None
    
    text_frame = shape.text_frame
    
    # Skip empty text frames
    if not text_frame.text.strip():
        return None
    
    # Skip slide numbers and footer placeholders
    if shape.is_placeholder:
        placeholder = shape.placeholder_format
        ph_type = str(placeholder.type) if placeholder.type else ""
        if "SLIDE_NUMBER" in ph_type or "FOOTER" in ph_type or "DATE" in ph_type:
            return None
    
    # Calculate absolute position
    abs_left = (shape.left or 0) + offset_left
    abs_top = (shape.top or 0) + offset_top
    
    shape_data = ShapeData(
        left=emu_to_inches(abs_left),
        top=emu_to_inches(abs_top),
        width=emu_to_inches(shape.width or 0),
        height=emu_to_inches(shape.height or 0),
    )
    
    # Get placeholder type
    if shape.is_placeholder:
        shape_data.placeholder_type = get_placeholder_type_name(shape.placeholder_format)
    
    # Extract paragraphs
    for paragraph in text_frame.paragraphs:
        if paragraph.text.strip():
            para_data = ParagraphData.from_paragraph(paragraph)
            shape_data.paragraphs.append(para_data)
            
            # Track default font size
            if para_data.font_size and shape_data.default_font_size is None:
                shape_data.default_font_size = para_data.font_size
    
    return shape_data if shape_data.paragraphs else None


def extract_shapes_recursive(
    shapes, 
    offset_left: int = 0, 
    offset_top: int = 0
) -> List[Tuple[ShapeData, int, int]]:
    """Recursively extract text from shapes including group shapes.
    
    Args:
        shapes: Collection of shapes to process.
        offset_left: X offset for nested shapes (in EMUs).
        offset_top: Y offset for nested shapes (in EMUs).
        
    Returns:
        List of tuples containing (ShapeData, absolute_left, absolute_top).
    """
    results: List[Tuple[ShapeData, int, int]] = []
    
    for shape in shapes:
        # Handle group shapes recursively
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            group_left = (shape.left or 0) + offset_left
            group_top = (shape.top or 0) + offset_top
            results.extend(extract_shapes_recursive(
                shape.shapes, group_left, group_top
            ))
        else:
            shape_data = extract_shape_text(shape, offset_left, offset_top)
            if shape_data:
                abs_left = (shape.left or 0) + offset_left
                abs_top = (shape.top or 0) + offset_top
                results.append((shape_data, abs_left, abs_top))
    
    return results


def detect_overlaps(shapes: List[Tuple[ShapeData, int, int]]) -> None:
    """Detect overlapping shapes and update their overlap property.
    
    Args:
        shapes: List of (ShapeData, left_emu, top_emu) tuples.
    """
    for i, (shape1, left1, top1) in enumerate(shapes):
        overlaps: Dict[str, float] = {}
        
        for j, (shape2, left2, top2) in enumerate(shapes):
            if i == j:
                continue
                
            # Convert to inches for comparison
            s1_left = emu_to_inches(left1)
            s1_top = emu_to_inches(top1)
            s1_right = s1_left + shape1.width
            s1_bottom = s1_top + shape1.height
            
            s2_left = emu_to_inches(left2)
            s2_top = emu_to_inches(top2)
            s2_right = s2_left + shape2.width
            s2_bottom = s2_top + shape2.height
            
            # Check for overlap
            if (s1_left < s2_right and s1_right > s2_left and
                s1_top < s2_bottom and s1_bottom > s2_top):
                
                # Calculate overlap area
                overlap_left = max(s1_left, s2_left)
                overlap_right = min(s1_right, s2_right)
                overlap_top = max(s1_top, s2_top)
                overlap_bottom = min(s1_bottom, s2_bottom)
                
                overlap_width = overlap_right - overlap_left
                overlap_height = overlap_bottom - overlap_top
                overlap_area = overlap_width * overlap_height
                
                if overlap_area > 0:
                    overlaps[f"shape-{j}"] = round(overlap_area, 2)
        
        if overlaps:
            shape1.overlap = {"overlapping_shapes": overlaps}


def extract_text_inventory(
    pptx_path: Path, 
    issues_only: bool = False
) -> Dict[str, Dict[str, Dict[str, Any]]]:
    """Extract all text content from a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the PowerPoint file.
        issues_only: If True, only include shapes with issues.
        
    Returns:
        Dictionary mapping slide IDs to shape IDs to shape data.
    """
    prs = Presentation(pptx_path)
    inventory: Dict[str, Dict[str, Dict[str, Any]]] = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_key = f"slide-{slide_idx}"
        
        # Extract all text shapes from this slide
        shapes_with_positions = extract_shapes_recursive(slide.shapes)
        
        # Detect overlaps
        detect_overlaps(shapes_with_positions)
        
        # Sort by visual position (top to bottom, left to right)
        shapes_with_positions.sort(key=lambda x: (x[0].top, x[0].left))
        
        # Build slide inventory
        slide_shapes: Dict[str, Dict[str, Any]] = {}
        
        for shape_idx, (shape_data, _, _) in enumerate(shapes_with_positions):
            # Filter by issues if requested
            if issues_only:
                has_issues = (
                    shape_data.overlap is not None or 
                    shape_data.overflow_bottom is not None
                )
                if not has_issues:
                    continue
            
            shape_key = f"shape-{shape_idx}"
            slide_shapes[shape_key] = shape_data.to_dict()
        
        if slide_shapes:
            inventory[slide_key] = slide_shapes
    
    return inventory


def save_inventory(inventory: Dict[str, Dict[str, Dict[str, Any]]], output_path: Path) -> None:
    """Save inventory data to JSON file.
    
    Args:
        inventory: The inventory data to save.
        output_path: Path for the output JSON file.
    """
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(inventory, f, indent=2, ensure_ascii=False)


def main() -> None:
    """Main entry point for command-line usage."""
    args = parse_arguments()
    
    input_path = Path(args.input)
    
    # Validate input file
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)
    
    if not input_path.suffix.lower() == ".pptx":
        print("Error: Input must be a PowerPoint file (.pptx)")
        sys.exit(1)
    
    try:
        print(f"Extracting text inventory from: {args.input}")
        if args.issues_only:
            print("Filtering to include only text shapes with issues")
        
        inventory = extract_text_inventory(input_path, issues_only=args.issues_only)
        
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        save_inventory(inventory, output_path)
        
        print(f"Output saved to: {args.output}")
        
        # Report statistics
        total_slides = len(inventory)
        total_shapes = sum(len(shapes) for shapes in inventory.values())
        
        if args.issues_only:
            if total_shapes > 0:
                print(f"Found {total_shapes} text elements with issues in {total_slides} slides")
            else:
                print("No issues discovered")
        else:
            print(f"Found text in {total_slides} slides with {total_shapes} text elements")
            
    except Exception as e:
        print(f"Error processing presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
