# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
#
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
insert_diagram_slides.py - 図形スライドを正しい位置・正しいレイアウトで挿入

Usage:
    python scripts/insert_diagram_slides.py <base.pptx> <diagrams.pptx> <output.pptx> --config <config.json>

config.json の形式:
{
  "insertions": [
    {"source_index": 0, "target_position": 4, "layout_name": "タイトルとコンテンツ"},
    {"source_index": 1, "target_position": 7, "layout_name": "タイトルとコンテンツ"},
    {"source_index": 2, "target_position": 10, "layout_name": "タイトルとコンテンツ"}
  ]
}
"""

import argparse
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu


def find_layout_by_name(prs, layout_name: str):
    """レイアウト名でレイアウトを検索"""
    for slide_master in prs.slide_masters:
        for layout in slide_master.slide_layouts:
            if layout.name == layout_name:
                return layout
    return None


def find_content_layout(prs):
    """コンテンツ系レイアウトを検索（フォールバック用）"""
    keywords = ["コンテンツ", "Content", "content", "タイトルとコンテンツ"]
    for slide_master in prs.slide_masters:
        for layout in slide_master.slide_layouts:
            for kw in keywords:
                if kw.lower() in layout.name.lower():
                    return layout
    # 見つからない場合は最初のレイアウト
    return prs.slide_masters[0].slide_layouts[1] if len(prs.slide_masters[0].slide_layouts) > 1 else prs.slide_masters[0].slide_layouts[0]


def copy_shapes_to_slide(source_slide, target_slide, scale_x=1.0, scale_y=1.0, offset_x=0, offset_y=0):
    """ソーススライドのシェイプをターゲットスライドにコピー"""
    copied = 0
    for shape in source_slide.shapes:
        try:
            sp = shape._element
            new_sp = copy.deepcopy(sp)
            
            # スケーリングとオフセットを適用
            if scale_x != 1.0 or scale_y != 1.0 or offset_x != 0 or offset_y != 0:
                # xfrm 要素を取得
                xfrm = new_sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                if xfrm is not None:
                    off = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                    ext = xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                    if off is not None:
                        x = int(off.get('x', 0))
                        y = int(off.get('y', 0))
                        off.set('x', str(int(x * scale_x + offset_x)))
                        off.set('y', str(int(y * scale_y + offset_y)))
                    if ext is not None:
                        cx = int(ext.get('cx', 0))
                        cy = int(ext.get('cy', 0))
                        ext.set('cx', str(int(cx * scale_x)))
                        ext.set('cy', str(int(cy * scale_y)))
            
            target_slide.shapes._spTree.append(new_sp)
            copied += 1
        except Exception as e:
            print(f"    [!] Shape copy failed: {e}")
    
    return copied


def insert_diagram_slides(base_path: str, diagrams_path: str, output_path: str, config: dict):
    """
    図形スライドを正しい位置に挿入
    
    Args:
        base_path: ベースPPTX（テンプレートから生成したもの）
        diagrams_path: 図形スライドPPTX（pptxgenjsで生成）
        output_path: 出力先
        config: 挿入設定
    """
    print(f"\n📋 Insert Diagram Slides")
    print(f"  Base:     {base_path}")
    print(f"  Diagrams: {diagrams_path}")
    print(f"  Output:   {output_path}")
    print()
    
    # PPTXを読み込み
    base_prs = Presentation(base_path)
    diagrams_prs = Presentation(diagrams_path)
    
    # サイズを取得
    base_w = base_prs.slide_width.emu
    base_h = base_prs.slide_height.emu
    diag_w = diagrams_prs.slide_width.emu
    diag_h = diagrams_prs.slide_height.emu
    
    base_w_in = base_prs.slide_width.inches
    diag_w_in = diagrams_prs.slide_width.inches
    
    print(f"  Base size:     {base_w_in:.2f}\" x {base_prs.slide_height.inches:.2f}\"")
    print(f"  Diagrams size: {diag_w_in:.2f}\" x {diagrams_prs.slide_height.inches:.2f}\"")
    
    # スケール計算
    scale_x = base_w / diag_w if diag_w != base_w else 1.0
    scale_y = base_h / diag_h if diag_h != base_h else 1.0
    
    if scale_x != 1.0 or scale_y != 1.0:
        print(f"  Scale: {scale_x:.2%} x {scale_y:.2%}")
    print()
    
    # 挿入設定を位置でソート（後ろから挿入するため降順）
    insertions = sorted(config.get("insertions", []), key=lambda x: x["target_position"], reverse=True)
    
    print(f"  Insertions: {len(insertions)}")
    for ins in reversed(insertions):  # 表示は昇順
        print(f"    - Source[{ins['source_index']}] → Position {ins['target_position']} ({ins.get('layout_name', 'default')})")
    print()
    
    # 後ろから挿入（インデックスがずれないように）
    for ins in insertions:
        source_idx = ins["source_index"]
        target_pos = ins["target_position"]
        layout_name = ins.get("layout_name", "タイトルとコンテンツ")
        
        if source_idx >= len(diagrams_prs.slides):
            print(f"  [!] Source index {source_idx} out of range, skipping")
            continue
        
        source_slide = diagrams_prs.slides[source_idx]
        
        # レイアウトを取得
        layout = find_layout_by_name(base_prs, layout_name)
        if layout is None:
            print(f"  [!] Layout '{layout_name}' not found, using fallback")
            layout = find_content_layout(base_prs)
        
        # 新しいスライドを挿入
        # python-pptx では add_slide は末尾に追加するので、XML操作で位置を調整
        new_slide = base_prs.slides.add_slide(layout)
        
        # スライドの位置を調整
        slide_id = new_slide.slide_id
        sldIdLst = base_prs.part._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
        
        if sldIdLst is not None:
            # 最後に追加されたスライドを取得
            sld_ids = list(sldIdLst)
            if len(sld_ids) > 0:
                last_sld_id = sld_ids[-1]
                # 目的の位置に移動
                if target_pos < len(sld_ids) - 1:
                    target_sld_id = sld_ids[target_pos]
                    sldIdLst.remove(last_sld_id)
                    sldIdLst.insert(target_pos, last_sld_id)
        
        # プレースホルダーを削除（図形で上書きするため）
        for shape in list(new_slide.shapes):
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                pass
        
        # シェイプをコピー
        copied = copy_shapes_to_slide(source_slide, new_slide, scale_x, scale_y)
        
        # ノートをコピー
        if source_slide.has_notes_slide:
            try:
                notes_text = source_slide.notes_slide.notes_text_frame.text
                if notes_text:
                    if not new_slide.has_notes_slide:
                        new_slide.notes_slide
                    new_slide.notes_slide.notes_text_frame.text = notes_text
            except:
                pass
        
        print(f"  [Slide {target_pos + 1}] Inserted with {copied} shapes (Layout: {layout.name})")
    
    # 保存
    base_prs.save(output_path)
    print()
    print(f"✅ Created: {output_path}")
    print(f"   Total slides: {len(base_prs.slides)}")
    
    return len(base_prs.slides)


def main():
    parser = argparse.ArgumentParser(description="Insert diagram slides at specific positions")
    parser.add_argument("base", help="Base PPTX (from template)")
    parser.add_argument("diagrams", help="Diagram slides PPTX (from pptxgenjs)")
    parser.add_argument("output", help="Output PPTX")
    parser.add_argument("--config", required=True, help="Config JSON file")
    
    args = parser.parse_args()
    
    # 設定を読み込み
    with open(args.config, "r", encoding="utf-8") as f:
        config = json.load(f)
    
    insert_diagram_slides(args.base, args.diagrams, args.output, config)


if __name__ == "__main__":
    main()
