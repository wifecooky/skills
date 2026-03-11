# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
#
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
merge_slides.py - pptxgenjs で生成した構成図をテンプレートにマージ

Usage:
    python scripts/merge_slides.py <template.pptx> <source.pptx> <output.pptx>
    python scripts/merge_slides.py <template.pptx> <source.pptx> <output.pptx> --position 0

Arguments:
    template.pptx  : マスター継承元のテンプレート
    source.pptx    : pptxgenjs 等で生成した構成図スライド
    output.pptx    : 出力先

Options:
    --position N   : 挿入位置（0=先頭, -1=末尾, default: -1）
    --clear-template : テンプレートの既存スライドを削除（マスターのみ継承）
    --keep-source-master : ソースのマスター色を維持（テンプレート色に変換しない）

Examples:
    # 構成図をテンプレートの末尾にマージ
    python scripts/merge_slides.py templates/Mytemplate_MS.pptx output_ppt/diagram.pptx output_ppt/merged.pptx

    # テンプレートのスライドを削除してマスターのみ継承（★ 推奨）
    python scripts/merge_slides.py templates/Mytemplate_MS.pptx output_ppt/diagram.pptx output_ppt/merged.pptx --clear-template

    # 先頭に挿入
    python scripts/merge_slides.py templates/Mytemplate_MS.pptx output_ppt/diagram.pptx output_ppt/merged.pptx --position 0
"""

import argparse
import copy
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree


def copy_shape(shape, target_slide):
    """
    シェイプを別のスライドにコピーする
    
    Args:
        shape: コピー元のシェイプ
        target_slide: コピー先のスライド
    
    Returns:
        コピーされたシェイプ（または None）
    """
    try:
        # XML要素を直接コピー（最も確実な方法）
        sp = shape._element
        new_sp = copy.deepcopy(sp)
        target_slide.shapes._spTree.append(new_sp)
        return new_sp
    except Exception as e:
        print(f"  [!] Failed to copy shape: {shape.shape_type} - {e}")
        return None


def merge_slides(
    template_path: str,
    source_path: str,
    output_path: str,
    position: int = -1,
    keep_source_master: bool = False,
    clear_template: bool = False
) -> int:
    """
    テンプレートに pptxgenjs で作成したスライドをマージする
    
    Args:
        template_path: テンプレート PPTX（マスター継承元）
        source_path: pptxgenjs で生成した構成図 PPTX
        output_path: 出力先
        position: 挿入位置（-1=末尾, 0=先頭, N=N番目の後）
        keep_source_master: True の場合、ソースの色を維持
        clear_template: True の場合、テンプレートの既存スライドを削除
    
    Returns:
        マージされたスライド数
    """
    print(f"\n📋 Merge Slides")
    print(f"  Template: {template_path}")
    print(f"  Source:   {source_path}")
    print(f"  Output:   {output_path}")
    print(f"  Position: {'末尾' if position == -1 else f'{position}番目'}")
    
    # ファイル存在確認
    if not Path(template_path).exists():
        print(f"❌ Template not found: {template_path}")
        return 0
    if not Path(source_path).exists():
        print(f"❌ Source not found: {source_path}")
        return 0
    
    # プレゼンテーションを開く
    template = Presentation(template_path)
    source = Presentation(source_path)
    
    print(f"\n  Template slides: {len(template.slides)}")
    print(f"  Source slides:   {len(source.slides)}")
    
    # テンプレートのスライドを削除（--clear-template オプション）
    if clear_template and len(template.slides) > 0:
        print(f"\n  🗑️ Clearing {len(template.slides)} template slides (keeping master only)...")
        # スライドをすべて削除（逆順で削除）
        xml_slides = template.slides._sldIdLst
        slides_to_remove = list(xml_slides)
        for sldId in slides_to_remove:
            rId = sldId.rId
            template.part.drop_rel(rId)
            xml_slides.remove(sldId)
        print(f"  ✅ Template slides cleared")
    
    # テンプレートのスライドサイズを取得
    template_width = template.slide_width
    template_height = template.slide_height
    source_width = source.slide_width
    source_height = source.slide_height
    
    # サイズ比率を計算（スケーリング用）
    width_ratio = template_width / source_width
    height_ratio = template_height / source_height
    scale_ratio = min(width_ratio, height_ratio)  # アスペクト比維持
    
    if abs(scale_ratio - 1.0) > 0.01:
        print(f"\n  ⚠️ Size difference detected:")
        print(f"     Template: {template_width.inches:.2f}\" x {template_height.inches:.2f}\"")
        print(f"     Source:   {source_width.inches:.2f}\" x {source_height.inches:.2f}\"")
        print(f"     Scale:    {scale_ratio:.2%}")
    
    # Blank レイアウトを探す（最後のレイアウトを使用）
    blank_layout = None
    for layout in template.slide_layouts:
        if "blank" in layout.name.lower() or layout.name == "Blank":
            blank_layout = layout
            break
    
    if blank_layout is None:
        # Blank が見つからない場合は最後のレイアウトを使用
        blank_layout = template.slide_layouts[-1]
        print(f"  Using layout: {blank_layout.name} (fallback)")
    else:
        print(f"  Using layout: {blank_layout.name}")
    
    # 挿入位置を計算
    if position == -1:
        insert_index = len(template.slides)
    elif position >= 0:
        insert_index = min(position, len(template.slides))
    else:
        insert_index = len(template.slides)
    
    merged_count = 0
    
    # 各スライドをコピー
    for slide_idx, slide in enumerate(source.slides):
        print(f"\n  [Slide {slide_idx + 1}] Copying {len(slide.shapes)} shapes...")
        
        # 新しいスライドを追加
        new_slide = template.slides.add_slide(blank_layout)
        
        # スライドを正しい位置に移動
        if insert_index < len(template.slides) - 1:
            # スライドの順序を変更
            slide_id = new_slide.slide_id
            # XMLで順序を操作
            slides_part = template.slides._sldIdLst
            slide_elem = slides_part[-1]  # 追加したスライド
            slides_part.remove(slide_elem)
            slides_part.insert(insert_index + merged_count, slide_elem)
        
        # シェイプをコピー
        shape_count = 0
        for shape in slide.shapes:
            # プレースホルダーはスキップ（空の場合）
            if shape.is_placeholder:
                continue
            
            result = copy_shape(shape, new_slide)
            if result is not None:
                shape_count += 1
                
                # スケーリングが必要な場合
                if abs(scale_ratio - 1.0) > 0.01:
                    try:
                        # 位置とサイズをスケーリング
                        if hasattr(shape, 'left'):
                            result.set(qn('p:x'), str(int(int(result.get(qn('p:x'), 0)) * scale_ratio)))
                        if hasattr(shape, 'top'):
                            result.set(qn('p:y'), str(int(int(result.get(qn('p:y'), 0)) * scale_ratio)))
                    except:
                        pass
        
        # ノートをコピー
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text:
                notes_slide = new_slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
                print(f"    📝 Notes copied")
        
        print(f"    ✅ Copied {shape_count} shapes")
        merged_count += 1
    
    # 保存
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    template.save(output_path)
    
    print(f"\n✅ Merged {merged_count} slides")
    print(f"   Output: {output_path}")
    print(f"   Total slides: {len(template.slides)}")
    
    return merged_count


def main():
    parser = argparse.ArgumentParser(
        description="pptxgenjs で生成した構成図をテンプレートにマージ"
    )
    parser.add_argument("template", help="テンプレート PPTX（マスター継承元）")
    parser.add_argument("source", help="pptxgenjs 等で生成した構成図 PPTX")
    parser.add_argument("output", help="出力先 PPTX")
    parser.add_argument(
        "--position", type=int, default=-1,
        help="挿入位置（0=先頭, -1=末尾, N=N番目の後）"
    )
    parser.add_argument(
        "--clear-template", action="store_true",
        help="テンプレートの既存スライドを削除（マスターのみ継承）"
    )
    parser.add_argument(
        "--keep-source-master", action="store_true",
        help="ソースのマスター色を維持"
    )
    
    args = parser.parse_args()
    
    count = merge_slides(
        args.template,
        args.source,
        args.output,
        args.position,
        args.keep_source_master,
        getattr(args, 'clear_template', False)
    )
    
    sys.exit(0 if count > 0 else 1)


if __name__ == "__main__":
    main()
