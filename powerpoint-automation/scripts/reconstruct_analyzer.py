# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Analyze English PPTX content and convert to IR (content.json) for reconstruction.

This script takes an English PPTX file, extracts its content using extract_shapes.py,
analyzes the slide structure, and generates a content.json that can be used with
create_from_template.py to reconstruct the presentation using the original's
slide master layouts.

Key features:
- Analyzes slide content to determine optimal layout type
- Preserves text hierarchy (title, body, bullets)
- Identifies images for extraction
- Maps to appropriate slide master layouts

Usage:
    python scripts/reconstruct_analyzer.py <input.pptx> <output_content.json> [--layouts <layouts.json>]
    python scripts/reconstruct_analyzer.py <input.pptx> --analyze-only

Examples:
    # Full pipeline: analyze and generate content.json
    python scripts/reconstruct_analyzer.py input/presentation.pptx output_manifest/content.json

    # With pre-analyzed layouts
    python scripts/reconstruct_analyzer.py input/presentation.pptx output_manifest/content.json --layouts output_manifest/custom_layouts.json

    # Just analyze structure without generating content.json
    python scripts/reconstruct_analyzer.py input/presentation.pptx --analyze-only
"""

import argparse
import json
import sys
import io
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from dataclasses import dataclass, field

# Fix Windows console encoding issues
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class SlideContent:
    """Represents extracted content from a single slide."""
    slide_index: int
    title: str = ""
    subtitle: str = ""
    body_items: List[str] = field(default_factory=list)
    bullet_levels: List[int] = field(default_factory=list)
    has_image: bool = False
    image_count: int = 0
    has_table: bool = False
    has_chart: bool = False
    has_diagram: bool = False
    shape_count: int = 0
    text_density: str = "normal"  # low, normal, high
    layout_name: str = ""
    detected_type: str = "content"  # title, section, content, two_column, photo, closing, etc.
    notes: str = ""


def analyze_slide_layout(slide) -> Tuple[str, List[Dict]]:
    """Analyze the slide's current layout and extract placeholder info.
    
    Args:
        slide: PowerPoint slide object.
        
    Returns:
        Tuple of (layout_name, list of placeholder info dicts).
    """
    layout_name = slide.slide_layout.name or "Unknown"
    placeholders = []
    
    for shape in slide.slide_layout.placeholders:
        ph_type = str(shape.placeholder_format.type).replace("PLACEHOLDER_TYPE.", "")
        placeholders.append({
            "idx": shape.placeholder_format.idx,
            "type": ph_type,
        })
    
    return layout_name, placeholders


def extract_title_and_body(slide) -> Tuple[str, str, List[str], List[int]]:
    """Extract title, subtitle, and body items from a slide.
    
    Args:
        slide: PowerPoint slide object.
        
    Returns:
        Tuple of (title, subtitle, body_items, bullet_levels).
    """
    title = ""
    subtitle = ""
    body_items = []
    bullet_levels = []
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # Check if this is a placeholder
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type)
            
            # Handle title placeholders
            if "TITLE" in ph_type or "CENTER_TITLE" in ph_type:
                text = shape.text_frame.text.strip()
                # Clean soft line breaks
                text = text.replace('\x0b', ' ').replace('\u000b', ' ')
                title = text
                continue
            
            # Handle subtitle
            if "SUBTITLE" in ph_type:
                text = shape.text_frame.text.strip()
                text = text.replace('\x0b', ' ').replace('\u000b', ' ')
                subtitle = text
                continue
            
            # Handle body/content placeholders
            if "BODY" in ph_type or "OBJECT" in ph_type or "CONTENT" in ph_type:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text = text.replace('\x0b', ' ').replace('\u000b', ' ')
                        body_items.append(text)
                        
                        # Detect bullet level from paragraph properties
                        level = paragraph.level if paragraph.level else 0
                        bullet_levels.append(level)
                continue
        
        # Non-placeholder text shapes - add to body if substantial
        if shape.text_frame.text.strip():
            # Skip slide numbers, footers, etc. (typically small shapes at bottom)
            if shape.top and shape.height:
                # Check if it's in the footer area (bottom 15% of slide)
                # Standard slide height is ~5.625 inches (16:9) = ~5143500 EMU
                SLIDE_HEIGHT_EMU = 5143500
                if shape.top > SLIDE_HEIGHT_EMU * 0.85:
                    continue
            
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and len(text) > 3:  # Skip very short text
                    text = text.replace('\x0b', ' ').replace('\u000b', ' ')
                    # Only add if not already captured
                    if text not in body_items and text != title and text != subtitle:
                        body_items.append(text)
                        bullet_levels.append(0)
    
    return title, subtitle, body_items, bullet_levels


def count_visual_elements(slide) -> Dict[str, int]:
    """Count images, tables, charts, and diagrams in a slide.
    
    Args:
        slide: PowerPoint slide object.
        
    Returns:
        Dict with counts for each element type.
    """
    counts = {
        "images": 0,
        "tables": 0,
        "charts": 0,
        "diagrams": 0,  # SmartArt, grouped shapes
        "shapes": 0,
    }
    
    def count_recursive(shapes):
        for shape in shapes:
            # Group shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                counts["diagrams"] += 1
                # Don't recurse into groups - treat as single diagram
                continue
            
            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["images"] += 1
                continue
            
            # Tables
            if shape.has_table:
                counts["tables"] += 1
                continue
            
            # Charts
            if shape.has_chart:
                counts["charts"] += 1
                continue
            
            # SmartArt / Diagrams
            if hasattr(shape, 'ole_format') or 'SmartArt' in str(type(shape)):
                counts["diagrams"] += 1
                continue
            
            # Other shapes (excluding placeholders and text boxes)
            if not shape.is_placeholder and not shape.has_text_frame:
                if shape.shape_type not in [MSO_SHAPE_TYPE.TEXT_BOX]:
                    counts["shapes"] += 1
    
    count_recursive(slide.shapes)
    return counts


def infer_title_from_notes(notes: str) -> str:
    """Infer a slide title from speaker notes when title is missing.
    
    Args:
        notes: Speaker notes text.
        
    Returns:
        Inferred title or empty string.
    """
    if not notes:
        return ""
    
    # Take first sentence or first line (up to 60 chars)
    lines = notes.strip().split('\n')
    first_line = lines[0].strip()
    
    # If first line is a question or statement, use it
    if first_line and len(first_line) <= 80:
        # Clean up common patterns
        if first_line.endswith((':',)):
            first_line = first_line[:-1]
        return first_line[:60] + ('...' if len(first_line) > 60 else '')
    
    # Try to extract key phrase from notes
    sentences = notes.replace('\n', ' ').split('.')
    if sentences and len(sentences[0].strip()) <= 60:
        return sentences[0].strip()
    
    return ""


def is_empty_slide(content: SlideContent) -> tuple[bool, str]:
    """Check if a slide is empty (only notes or background).
    
    Args:
        content: SlideContent object with extracted data.
        
    Returns:
        Tuple of (is_empty, reason).
    """
    has_title = bool(content.title and content.title.strip())
    has_body = bool(content.body_items)
    has_visual = content.has_image or content.has_chart or content.has_table or content.has_diagram
    has_notes = bool(content.notes and content.notes.strip())
    
    # Slide with ONLY speaker notes = empty (useless for presentation)
    if has_notes and not has_title and not has_body and not has_visual:
        return True, "notes_only"
    
    # Completely blank slide (no content at all)
    if not has_title and not has_body and not has_visual and not has_notes:
        return True, "blank"
    
    # Title-only slides are OK (section headers, etc.)
    return False, ""


def detect_slide_type(content: SlideContent, layout_name: str) -> str:
    """Detect the appropriate slide type based on content analysis.
    
    Args:
        content: SlideContent object with extracted data.
        layout_name: Name of the current layout.
        
    Returns:
        Detected slide type string.
    """
    layout_lower = layout_name.lower()
    title_lower = content.title.lower() if content.title else ""
    
    # Title slide detection
    if "title slide" in layout_lower or "タイトル スライド" in layout_lower:
        return "title"
    if content.subtitle and not content.body_items and not content.has_image:
        return "title"
    
    # Section header detection
    if "section" in layout_lower or "セクション" in layout_lower:
        return "section"
    if len(content.body_items) == 0 and content.title and not content.subtitle:
        # Just a title, likely a section divider
        if len(content.title) < 50:  # Short title
            return "section"
    
    # Closing slide detection
    if "closing" in layout_lower or "end" in layout_lower:
        return "closing"
    closing_keywords = ["thank you", "thanks", "questions", "q&a", "contact", 
                       "ありがとう", "質問", "お問い合わせ"]
    if any(kw in title_lower for kw in closing_keywords):
        return "closing"
    
    # Agenda detection
    agenda_keywords = ["agenda", "outline", "contents", "目次", "アジェンダ", "概要"]
    if any(kw in title_lower for kw in agenda_keywords):
        return "agenda"
    
    # Photo/image-heavy slide
    if content.has_image and content.image_count >= 1:
        if len(content.body_items) <= 2:
            return "photo"
    
    # Two-column detection (based on layout name or content pattern)
    if "two column" in layout_lower or "2列" in layout_lower or "comparison" in layout_lower:
        return "two_column"
    
    # Summary detection
    summary_keywords = ["summary", "まとめ", "conclusion", "key takeaway", "結論"]
    if any(kw in title_lower for kw in summary_keywords):
        return "summary"
    
    # Default to content
    return "content"


def extract_slide_notes(slide) -> str:
    """Extract speaker notes from a slide.
    
    Args:
        slide: PowerPoint slide object.
        
    Returns:
        Notes text or empty string.
    """
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            return notes_frame.text.strip()
    return ""


def analyze_presentation(pptx_path: Path) -> List[SlideContent]:
    """Analyze all slides in a presentation.
    
    Args:
        pptx_path: Path to the PowerPoint file.
        
    Returns:
        List of SlideContent objects.
    """
    prs = Presentation(pptx_path)
    slides_content = []
    
    for idx, slide in enumerate(prs.slides):
        # Get layout info
        layout_name, placeholders = analyze_slide_layout(slide)
        
        # Extract text content
        title, subtitle, body_items, bullet_levels = extract_title_and_body(slide)
        
        # Count visual elements
        counts = count_visual_elements(slide)
        
        # Get notes
        notes = extract_slide_notes(slide)
        
        # Create content object
        content = SlideContent(
            slide_index=idx,
            title=title,
            subtitle=subtitle,
            body_items=body_items,
            bullet_levels=bullet_levels,
            has_image=counts["images"] > 0,
            image_count=counts["images"],
            has_table=counts["tables"] > 0,
            has_chart=counts["charts"] > 0,
            has_diagram=counts["diagrams"] > 0,
            shape_count=sum(counts.values()),
            layout_name=layout_name,
            notes=notes,
        )
        
        # Detect text density
        total_text = len(title) + len(subtitle) + sum(len(item) for item in body_items)
        if total_text < 50:
            content.text_density = "low"
        elif total_text > 300:
            content.text_density = "high"
        else:
            content.text_density = "normal"
        
        # Detect slide type
        content.detected_type = detect_slide_type(content, layout_name)
        
        # Check for empty slides (notes only, blank)
        is_empty, reason = is_empty_slide(content)
        if is_empty:
            if reason == "notes_only":
                print(f"  ⚠️  Slide {idx + 1}: EMPTY - only has speaker notes (no visible content)")
            elif reason == "blank":
                print(f"  ⚠️  Slide {idx + 1}: BLANK - no content at all")
            content.detected_type = "_empty"  # Mark for filtering
        
        # If title is empty, try to infer from notes
        if not content.title and content.notes:
            inferred = infer_title_from_notes(content.notes)
            if inferred:
                content.title = inferred
                print(f"  Slide {idx + 1}: Inferred title from notes: '{inferred[:40]}...'")
        
        slides_content.append(content)
    
    return slides_content


def generate_content_json(
    slides_content: List[SlideContent],
    layouts_mapping: Optional[Dict] = None
) -> Dict[str, Any]:
    """Generate content.json structure from analyzed slides.
    
    Args:
        slides_content: List of SlideContent objects.
        layouts_mapping: Optional layout mapping from analyze_template.py.
        
    Returns:
        Dict structure for content.json.
    """
    result = {
        "title": "",
        "subtitle": "",
        "slides": [],
    }
    
    # Extract presentation title from first title slide
    for content in slides_content:
        if content.detected_type == "title":
            result["title"] = content.title
            result["subtitle"] = content.subtitle
            break
    
    # Process each slide
    for content in slides_content:
        slide_data = {
            "type": content.detected_type,
            "title": content.title,
        }
        
        # Add subtitle for title/section slides
        if content.subtitle and content.detected_type in ["title", "section"]:
            slide_data["subtitle"] = content.subtitle
        
        # Add body items
        if content.body_items:
            # Check if we should use 'items' or 'bullets' format
            if content.detected_type in ["content", "agenda", "summary", "two_column"]:
                # Convert to items with bullet structure if needed
                if any(level > 0 for level in content.bullet_levels):
                    # Has nested bullets - use bullets format
                    slide_data["bullets"] = []
                    for item, level in zip(content.body_items, content.bullet_levels):
                        slide_data["bullets"].append({
                            "text": item,
                            "level": level,
                        })
                else:
                    # Flat list
                    slide_data["items"] = content.body_items
        
        # Add image placeholder if slide has images
        if content.has_image:
            slide_data["image"] = {
                "path": f"images/slide_{content.slide_index + 1:02d}.png",
                "position": "right",
                "width_percent": 45,
                "_note": "TODO: Extract actual image from original PPTX"
            }
        
        # Add notes if present, with source reference
        source_ref = f"[出典: 元スライド #{content.slide_index + 1}]"
        if content.notes:
            slide_data["notes"] = f"{content.notes}\n\n---\n{source_ref}"
        else:
            slide_data["notes"] = source_ref
        
        # Add metadata for debugging/review
        slide_data["_source"] = {
            "original_index": content.slide_index,
            "original_layout": content.layout_name,
            "has_chart": content.has_chart,
            "has_table": content.has_table,
            "has_diagram": content.has_diagram,
            "text_density": content.text_density,
        }
        
        result["slides"].append(slide_data)
    
    return result


def print_analysis_summary(slides_content: List[SlideContent]) -> None:
    """Print human-readable analysis summary.
    
    Args:
        slides_content: List of SlideContent objects.
    """
    print(f"\n{'='*60}")
    print(f"Slide Analysis Summary")
    print(f"{'='*60}\n")
    
    type_counts = {}
    for content in slides_content:
        t = content.detected_type
        type_counts[t] = type_counts.get(t, 0) + 1
    
    # Count empty slides
    empty_count = sum(1 for c in slides_content if c.detected_type == "_empty")
    
    print(f"Total slides: {len(slides_content)}")
    if empty_count > 0:
        print(f"⚠️  Empty slides (notes-only or blank): {empty_count}")
        print(f"   These should be EXCLUDED from output!")
    
    print(f"\nSlide types detected:")
    for t, count in sorted(type_counts.items()):
        if t == "_empty":
            print(f"  {t:15} : {count}  ⚠️  EXCLUDE THESE")
        else:
            print(f"  {t:15} : {count}")
    
    print(f"\n{'─'*60}")
    print(f"{'#':>3}  {'Type':12}  {'Layout':25}  Title")
    print(f"{'─'*60}")
    
    for content in slides_content:
        title_preview = content.title[:35] + "..." if len(content.title) > 35 else content.title
        flags = []
        if content.has_image:
            flags.append(f"🖼️{content.image_count}")
        if content.has_chart:
            flags.append("📊")
        if content.has_table:
            flags.append("📋")
        if content.has_diagram:
            flags.append("🔷")
        
        flag_str = " ".join(flags)
        print(f"{content.slide_index + 1:3}  {content.detected_type:12}  {content.layout_name[:25]:25}  {title_preview} {flag_str}")
    
    print(f"{'─'*60}\n")


def main():
    parser = argparse.ArgumentParser(
        description="Analyze English PPTX and generate content.json for reconstruction.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("output", nargs="?", help="Output content.json file")
    parser.add_argument("--layouts", help="Layout mapping JSON from analyze_template.py")
    parser.add_argument("--context", help="Context JSON from classify_input.py (for base name)")
    parser.add_argument("--analyze-only", action="store_true", 
                       help="Only analyze, don't generate content.json")
    parser.add_argument("--include-metadata", action="store_true",
                       help="Include _source metadata in output (for debugging)")
    
    args = parser.parse_args()
    
    input_path = Path(args.input)
    
    # Load context if provided
    context = None
    if args.context:
        try:
            with open(args.context, 'r', encoding='utf-8') as f:
                context = json.load(f)
                print(f"📋 Loaded context: {context.get('base_name', 'N/A')}")
        except Exception as e:
            print(f"Warning: Could not load context: {e}")
    
    # Validate input
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)
    
    if not input_path.suffix.lower() == ".pptx":
        print("Error: Input must be a PowerPoint file (.pptx)")
        sys.exit(1)
    
    # Analyze presentation
    print(f"Analyzing: {args.input}")
    slides_content = analyze_presentation(input_path)
    
    # Print summary
    print_analysis_summary(slides_content)
    
    if args.analyze_only:
        print("Analysis complete (--analyze-only mode)")
        return
    
    # Need output path
    if not args.output:
        # Auto-generate output path using context or input stem
        if context and context.get("base_name"):
            base_name = context["base_name"]
        else:
            base_name = input_path.stem
        output_path = Path("output_manifest") / f"{base_name}_content.json"
    else:
        output_path = Path(args.output)
    
    # Prevent output to input/ directory
    if output_path.parts and output_path.parts[0] == 'input':
        print("Error: Cannot output to input/ directory. Use output_manifest/ instead.")
        sys.exit(1)
    
    # Load layouts mapping if provided
    layouts_mapping = None
    if args.layouts:
        with open(args.layouts, 'r', encoding='utf-8') as f:
            layouts_data = json.load(f)
            layouts_mapping = layouts_data.get("layout_mapping", {})
    
    # Generate content.json
    content = generate_content_json(slides_content, layouts_mapping)
    
    # Remove metadata if not requested
    if not args.include_metadata:
        for slide in content["slides"]:
            if "_source" in slide:
                del slide["_source"]
            if "image" in slide and "_note" in slide["image"]:
                del slide["image"]["_note"]
    
    # Save output
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(content, f, ensure_ascii=False, indent=2)
    
    print(f"✅ Generated: {output_path}")
    print(f"\nNext steps:")
    print(f"  1. Review and edit {output_path}")
    print(f"  2. Translate content (Localizer agent or manual)")
    print(f"  3. Extract images if needed")
    print(f"  4. Run: python scripts/create_from_template.py {args.input} {output_path} output.pptx")


if __name__ == "__main__":
    main()
