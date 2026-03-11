#!/usr/bin/env python3
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Rearrange PowerPoint slides based on a sequence of indices.

This module provides functionality to reorder and duplicate slides in a
PowerPoint presentation. It uses python-pptx library to manipulate PPTX files.

Usage:
    python reorder_slides.py template.pptx output.pptx 0,3,3,5,2

This will create output.pptx using slides from template.pptx in the specified order.
Slides can be repeated (e.g., 3 appears twice).

Author: aktsmm
License: CC BY-NC 4.0
"""

import argparse
import shutil
import sys
from copy import deepcopy
from pathlib import Path
from typing import Dict, List, Set

from pptx import Presentation
from pptx.slide import Slide


def parse_arguments() -> argparse.Namespace:
    """Parse command line arguments.
    
    Returns:
        Parsed command line arguments.
    """
    parser = argparse.ArgumentParser(
        description="Rearrange PowerPoint slides based on a sequence of indices.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python reorder_slides.py template.pptx output.pptx 0,3,3,5,2
    Creates output.pptx using slides 0, 3 (twice), 5, and 2 from template.pptx

  python reorder_slides.py template.pptx output.pptx 5,3,1,2,4
    Creates output.pptx with slides reordered as specified

Note: Slide indices are 0-based (first slide is 0, second is 1, etc.)
        """,
    )
    parser.add_argument("template", help="Path to template PPTX file")
    parser.add_argument("output", help="Path for output PPTX file")
    parser.add_argument(
        "sequence", help="Comma-separated sequence of slide indices (0-based)"
    )
    return parser.parse_args()


def parse_slide_sequence(sequence_str: str) -> List[int]:
    """Parse comma-separated slide indices string.
    
    Args:
        sequence_str: Comma-separated string of slide indices.
        
    Returns:
        List of slide indices as integers.
        
    Raises:
        ValueError: If sequence format is invalid.
    """
    try:
        return [int(x.strip()) for x in sequence_str.split(",")]
    except ValueError as e:
        raise ValueError(
            "Invalid sequence format. Use comma-separated integers (e.g., 0,3,3,5,2)"
        ) from e


def validate_slide_indices(indices: List[int], total_slides: int) -> None:
    """Validate that all slide indices are within valid range.
    
    Args:
        indices: List of slide indices to validate.
        total_slides: Total number of slides in the presentation.
        
    Raises:
        ValueError: If any index is out of range.
    """
    for idx in indices:
        if idx < 0 or idx >= total_slides:
            raise ValueError(
                f"Slide index {idx} out of range (valid: 0-{total_slides - 1})"
            )


def copy_slide_elements(source_slide: Slide, target_slide: Slide) -> None:
    """Copy all elements from source slide to target slide.
    
    This function copies shapes including pictures, text boxes, and other elements
    while preserving their relationships (e.g., embedded images).
    
    Args:
        source_slide: The source slide to copy from.
        target_slide: The target slide to copy to.
    """
    # Collect image relationships from source slide
    image_rels: Dict[str, any] = {}
    for rel_id, rel in source_slide.part.rels.items():
        if "image" in rel.reltype or "media" in rel.reltype:
            image_rels[rel_id] = rel

    # Clear existing placeholder shapes from target
    for shape in list(target_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy all shapes from source to target
    for shape in source_slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Handle image references - update blip embed IDs
        blips = new_el.xpath(".//a:blip[@r:embed]")
        for blip in blips:
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                old_rel = image_rels[old_rId]
                new_rId = target_slide.part.rels.get_or_add(
                    old_rel.reltype, old_rel._target
                )
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    # Ensure all image relationships are added to target
    for rel_id, rel in image_rels.items():
        try:
            target_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            pass  # Relationship might already exist


def duplicate_slide(prs: Presentation, source_index: int) -> Slide:
    """Create a duplicate of the specified slide.
    
    Args:
        prs: The presentation object.
        source_index: Index of the slide to duplicate.
        
    Returns:
        The newly created duplicate slide.
    """
    source = prs.slides[source_index]
    
    # Create new slide with same layout as source
    new_slide = prs.slides.add_slide(source.slide_layout)
    
    # Copy all elements from source to new slide
    copy_slide_elements(source, new_slide)
    
    return new_slide


def delete_slide(prs: Presentation, index: int) -> None:
    """Delete a slide from the presentation.
    
    Args:
        prs: The presentation object.
        index: Index of the slide to delete.
    """
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def move_slide(prs: Presentation, from_index: int, to_index: int) -> None:
    """Move a slide from one position to another.
    
    Args:
        prs: The presentation object.
        from_index: Current index of the slide.
        to_index: Target index for the slide.
    """
    slides = prs.slides._sldIdLst
    slide_element = slides[from_index]
    slides.remove(slide_element)
    slides.insert(to_index, slide_element)


def rearrange_presentation(
    template_path: Path, 
    output_path: Path, 
    slide_sequence: List[int]
) -> None:
    """Create a new presentation with slides from template in specified order.
    
    Args:
        template_path: Path to the template PPTX file.
        output_path: Path for the output PPTX file.
        slide_sequence: List of slide indices (0-based) to include.
        
    Raises:
        ValueError: If any slide index is out of range.
    """
    # Copy template to output path to preserve dimensions and theme
    if template_path != output_path:
        shutil.copy2(template_path, output_path)
        prs = Presentation(output_path)
    else:
        prs = Presentation(template_path)

    total_slides = len(prs.slides)
    
    # Validate all indices
    validate_slide_indices(slide_sequence, total_slides)

    # Track slide positions and duplicates
    slide_map: List[int] = []
    duplicated: Dict[int, List[int]] = {}

    print(f"Processing {len(slide_sequence)} slides from template...")
    
    # Phase 1: Create duplicates for repeated slides
    for i, template_idx in enumerate(slide_sequence):
        count = slide_sequence.count(template_idx)
        
        if template_idx in duplicated and duplicated[template_idx]:
            # Use existing duplicate
            slide_map.append(duplicated[template_idx].pop(0))
            print(f"  [{i}] Using duplicate of slide {template_idx}")
        elif count > 1 and template_idx not in duplicated:
            # First occurrence of repeated slide - create all duplicates now
            slide_map.append(template_idx)
            duplicates = []
            duplicate_count = count - 1
            print(f"  [{i}] Using original slide {template_idx}, creating {duplicate_count} duplicate(s)")
            
            for _ in range(duplicate_count):
                duplicate_slide(prs, template_idx)
                duplicates.append(len(prs.slides) - 1)
            duplicated[template_idx] = duplicates
        else:
            # Unique slide
            slide_map.append(template_idx)
            print(f"  [{i}] Using original slide {template_idx}")

    # Phase 2: Delete unused slides (work backwards to maintain indices)
    slides_to_keep: Set[int] = set(slide_map)
    print(f"\nDeleting {len(prs.slides) - len(slides_to_keep)} unused slides...")
    
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in slides_to_keep:
            delete_slide(prs, i)
            # Update indices in slide_map after deletion
            slide_map = [idx - 1 if idx > i else idx for idx in slide_map]

    # Phase 3: Reorder slides to final sequence
    print(f"Reordering {len(slide_map)} slides to final sequence...")
    
    for target_pos in range(len(slide_map)):
        current_pos = slide_map[target_pos]
        if current_pos != target_pos:
            move_slide(prs, current_pos, target_pos)
            # Update slide_map to reflect the move
            for i in range(len(slide_map)):
                if slide_map[i] > current_pos and slide_map[i] <= target_pos:
                    slide_map[i] -= 1
                elif slide_map[i] < current_pos and slide_map[i] >= target_pos:
                    slide_map[i] += 1
            slide_map[target_pos] = target_pos

    # Save the result
    prs.save(output_path)
    print(f"\nSaved rearranged presentation to: {output_path}")
    print(f"Final presentation has {len(prs.slides)} slides")


def main() -> None:
    """Main entry point for command-line usage."""
    args = parse_arguments()

    # Parse slide sequence
    try:
        slide_sequence = parse_slide_sequence(args.sequence)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)

    # Validate template exists
    template_path = Path(args.template)
    if not template_path.exists():
        print(f"Error: Template file not found: {args.template}")
        sys.exit(1)

    # Create output directory if needed
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Execute rearrangement
    try:
        rearrange_presentation(template_path, output_path, slide_sequence)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error processing presentation: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
