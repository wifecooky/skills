# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
#
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
review_pptx.py - PPTX content extraction and review helper

Extract slide content from PPTX for AI review or quick inspection.
Outputs structured information including titles, body text, and notes.

Usage:
    python scripts/review_pptx.py <input.pptx> [--json] [--output <file>]

Options:
    --json      Output in JSON format for programmatic processing
    --output    Save output to file instead of stdout
    --summary   Show summary only (slide count, issues detected)
"""

from pptx import Presentation
from pptx.util import Inches
import sys
import argparse
import json
import re
from pathlib import Path


def extract_slide_content(slide, slide_num: int) -> dict:
    """Extract content from a single slide.
    
    Args:
        slide: python-pptx Slide object
        slide_num: 1-based slide number
        
    Returns:
        Dictionary with slide content
    """
    content = {
        "slide_number": slide_num,
        "title": None,
        "body_texts": [],
        "notes": None,
        "issues": []
    }
    
    # Extract title
    if slide.shapes.title and slide.shapes.title.text:
        content["title"] = slide.shapes.title.text.strip()
    else:
        content["issues"].append("no_title")
    
    # Extract body text
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            if shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    # Check if it's just a page number
                    if re.match(r'^\d{1,3}$', text):
                        content["issues"].append("page_number_only")
                    else:
                        content["body_texts"].append(text)
    
    if not content["body_texts"]:
        content["issues"].append("empty_body")
    
    # Extract notes
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            content["notes"] = notes_text
            # Check if notes is source-only
            if re.match(r'^\[出典:.*\]$', notes_text.strip()):
                content["issues"].append("source_only_notes")
        else:
            content["issues"].append("empty_notes")
    else:
        content["issues"].append("no_notes")
    
    return content


def review_pptx(path: str, output_format: str = "text") -> dict:
    """Review PPTX content and extract structured information.
    
    Args:
        path: Path to PPTX file
        output_format: "text" or "json"
        
    Returns:
        Dictionary with review results
    """
    prs = Presentation(path)
    
    result = {
        "file": path,
        "slide_count": len(prs.slides),
        "size": {
            "width": round(prs.slide_width.inches, 2),
            "height": round(prs.slide_height.inches, 2)
        },
        "slides": [],
        "summary": {
            "empty_slides": 0,
            "missing_notes": 0,
            "source_only_notes": 0,
            "page_number_only": 0
        }
    }
    
    for i, slide in enumerate(prs.slides):
        slide_content = extract_slide_content(slide, i + 1)
        result["slides"].append(slide_content)
        
        # Update summary
        if "empty_body" in slide_content["issues"]:
            result["summary"]["empty_slides"] += 1
        if "no_notes" in slide_content["issues"] or "empty_notes" in slide_content["issues"]:
            result["summary"]["missing_notes"] += 1
        if "source_only_notes" in slide_content["issues"]:
            result["summary"]["source_only_notes"] += 1
        if "page_number_only" in slide_content["issues"]:
            result["summary"]["page_number_only"] += 1
    
    return result


def print_text_report(result: dict):
    """Print human-readable text report."""
    print(f'ファイル: {result["file"]}')
    print(f'スライド数: {result["slide_count"]}')
    print(f'サイズ: {result["size"]["width"]} x {result["size"]["height"]} インチ')
    print()
    
    # Summary
    issues = []
    if result["summary"]["empty_slides"] > 0:
        issues.append(f'空スライド: {result["summary"]["empty_slides"]}枚')
    if result["summary"]["missing_notes"] > 0:
        issues.append(f'ノート欠落: {result["summary"]["missing_notes"]}枚')
    if result["summary"]["source_only_notes"] > 0:
        issues.append(f'出典のみノート: {result["summary"]["source_only_notes"]}枚')
    if result["summary"]["page_number_only"] > 0:
        issues.append(f'ページ番号のみ: {result["summary"]["page_number_only"]}枚')
    
    if issues:
        print("⚠️ 検出された問題:")
        for issue in issues:
            print(f"  - {issue}")
        print()
    else:
        print("✅ 問題なし")
        print()
    
    # Slide details
    for slide in result["slides"]:
        print(f'=== スライド {slide["slide_number"]} ===')
        
        # Title
        if slide["title"]:
            print(f'タイトル: {slide["title"][:100]}')
        else:
            print('タイトル: (なし)')
        
        # Body text
        if slide["body_texts"]:
            for text in slide["body_texts"][:5]:
                display = text[:150] + '...' if len(text) > 150 else text
                display = display.replace('\n', ' | ')
                print(f'  {display}')
        
        # Notes
        if slide["notes"]:
            notes_preview = slide["notes"][:100]
            notes_preview = notes_preview.replace('\n', ' ')
            print(f'  [ノート] {notes_preview}...')
        
        # Issues
        if slide["issues"]:
            issue_labels = {
                "no_title": "⚠️ タイトルなし",
                "empty_body": "⚠️ 本文なし",
                "page_number_only": "⚠️ ページ番号のみ",
                "no_notes": "⚠️ ノートなし",
                "empty_notes": "⚠️ ノート空",
                "source_only_notes": "⚠️ 出典のみ"
            }
            for issue in slide["issues"]:
                if issue in issue_labels:
                    print(f'  {issue_labels[issue]}')
        
        print()


def main():
    parser = argparse.ArgumentParser(
        description='Review PPTX content for AI review or inspection'
    )
    parser.add_argument('input', help='Input PPTX file path')
    parser.add_argument('--json', action='store_true', help='Output in JSON format')
    parser.add_argument('--output', '-o', help='Output file path')
    parser.add_argument('--summary', action='store_true', help='Show summary only')
    
    args = parser.parse_args()
    
    if not Path(args.input).exists():
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)
    
    result = review_pptx(args.input)
    
    if args.json:
        output = json.dumps(result, ensure_ascii=False, indent=2)
        if args.output:
            Path(args.output).write_text(output, encoding='utf-8')
            print(f"Output saved to: {args.output}")
        else:
            print(output)
    elif args.summary:
        print(f'ファイル: {result["file"]}')
        print(f'スライド数: {result["slide_count"]}')
        print(f'空スライド: {result["summary"]["empty_slides"]}')
        print(f'ノート欠落: {result["summary"]["missing_notes"]}')
        print(f'出典のみノート: {result["summary"]["source_only_notes"]}')
        
        # Exit code based on issues
        if result["summary"]["empty_slides"] > 0:
            sys.exit(1)
        elif result["summary"]["missing_notes"] > 3:
            sys.exit(2)
    else:
        if args.output:
            import io
            old_stdout = sys.stdout
            sys.stdout = io.StringIO()
            print_text_report(result)
            output = sys.stdout.getvalue()
            sys.stdout = old_stdout
            Path(args.output).write_text(output, encoding='utf-8')
            print(f"Output saved to: {args.output}")
        else:
            print_text_report(result)


if __name__ == '__main__':
    main()
