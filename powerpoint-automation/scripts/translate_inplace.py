#!/usr/bin/env python3
"""
translate_inplace.py - In-place PPTX translation preserving all formatting.

Unlike the content.json pipeline (reconstruct → translate → rebuild), this script
directly edits the original PPTX file, replacing text while preserving:
  - All layouts, images, charts, shapes, and positioning
  - Grouped/nested shapes (recursive traversal)
  - Run-level formatting (font, color, size, bold, italic)

Usage:
  python translate_inplace.py <input.pptx> <translations.json> <output.pptx>

  translations.json format:
  {
    "original text 1": "translated text 1",
    "original text 2": "translated text 2"
  }

  Or use --extract to dump all text for translation:
  python translate_inplace.py <input.pptx> --extract <output.json>

  Or use --auto with target language for AI-assisted extraction:
  python translate_inplace.py <input.pptx> --auto <output.json> --lang zh
"""

import argparse
import json
import sys
from pptx import Presentation


def extract_all_text(pptx_path):
    """Extract all translatable text from PPTX, including grouped shapes.

    Returns a dict mapping slide_number to list of texts found.
    Also returns a flat dict {text: text} for use as translation template.
    """
    prs = Presentation(pptx_path)
    by_slide = {}
    flat = {}

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        _collect_texts(slide.shapes, texts)
        by_slide[f"slide_{i}"] = texts
        for t in texts:
            if t not in flat:
                flat[t] = t  # identity mapping as template

    return by_slide, flat


def _collect_texts(shapes, out):
    """Recursively collect paragraph-level text from shapes."""
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    out.append(text)
        # Recurse into grouped shapes
        if hasattr(shape, 'shapes'):
            _collect_texts(shape.shapes, out)
        # Table cells
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            out.append(text)


def translate_pptx(pptx_path, translations, output_path):
    """Translate PPTX in-place using translation map.

    Args:
        pptx_path: Path to source PPTX
        translations: dict mapping original_text → translated_text
        output_path: Path for translated PPTX output
    """
    prs = Presentation(pptx_path)

    replaced = 0
    skipped = 0

    for slide in prs.slides:
        r, s = _translate_shapes(slide.shapes, translations)
        replaced += r
        skipped += s

    prs.save(output_path)
    print(f"\n✅ Saved: {output_path}")
    print(f"   Replaced: {replaced} paragraphs")
    print(f"   Skipped (no match): {skipped} paragraphs")
    return replaced


def _translate_shapes(shapes, translations):
    """Recursively translate text in shapes, preserving formatting."""
    replaced = 0
    skipped = 0

    for shape in shapes:
        if shape.has_text_frame:
            r, s = _translate_text_frame(shape.text_frame, translations)
            replaced += r
            skipped += s

        # Recurse into grouped shapes
        if hasattr(shape, 'shapes'):
            r, s = _translate_shapes(shape.shapes, translations)
            replaced += r
            skipped += s

        # Table cells
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    r, s = _translate_text_frame(cell.text_frame, translations)
                    replaced += r
                    skipped += s

    return replaced, skipped


def _translate_text_frame(text_frame, translations):
    """Translate paragraphs in a text frame, preserving run-level formatting.

    Strategy: Match full paragraph text against translation map.
    If matched, put translated text into the first run (preserving its formatting)
    and clear all subsequent runs. This preserves font/color/size of the first run.
    """
    replaced = 0
    skipped = 0

    for para in text_frame.paragraphs:
        full_text = para.text.strip()
        if not full_text:
            continue

        if full_text in translations:
            translated = translations[full_text]
            if translated == full_text:
                continue  # No change needed

            runs = para.runs
            if runs:
                runs[0].text = translated
                for r in runs[1:]:
                    r.text = ""
            replaced += 1
        else:
            skipped += 1

    return replaced, skipped


def main():
    parser = argparse.ArgumentParser(
        description='In-place PPTX translation preserving all formatting')
    parser.add_argument('input', help='Input PPTX file path')
    parser.add_argument('translations_or_output',
                        help='translations.json (translate mode) or output.json (extract mode)')
    parser.add_argument('output', nargs='?', help='Output PPTX file path')
    parser.add_argument('--extract', action='store_true',
                        help='Extract all text to JSON for translation')
    parser.add_argument('--lang', default=None,
                        help='Target language hint (used with --extract for comments)')

    args = parser.parse_args()

    if args.extract:
        # Extract mode
        by_slide, flat = extract_all_text(args.input)
        result = {
            "_meta": {
                "source": args.input,
                "target_lang": args.lang or "TODO",
                "total_strings": len(flat),
                "instruction": "Replace values with translated text. Keys are original text."
            },
            "by_slide": by_slide,
            "translations": flat
        }
        with open(args.translations_or_output, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"✅ Extracted {len(flat)} unique strings to: {args.translations_or_output}")
        print(f"   Edit the 'translations' section, then run:")
        print(f"   python translate_inplace.py {args.input} {args.translations_or_output} output.pptx")
        return

    # Translate mode
    if not args.output:
        print("Error: output PPTX path required in translate mode", file=sys.stderr)
        sys.exit(1)

    with open(args.translations_or_output, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Support both flat dict and nested format
    if 'translations' in data:
        translations = data['translations']
    else:
        translations = data

    print(f"Loading: {args.input}")
    print(f"Translations: {len(translations)} entries")
    translate_pptx(args.input, translations, args.output)


if __name__ == '__main__':
    main()
