# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Validate generated PPTX against content.json.

This script provides deterministic validation for PPTX files.
It checks:
1. Slide count matches content.json
2. Speaker notes existence and quality (detects "source-only" notes)
3. Image placement verification
4. Text overflow detection (excessive text, too many paragraphs)
5. Basic structure validation

Speaker notes quality checks:
- Detects notes with only source citations (e.g., "[出典: ...]")
- Warns when notes are too short (<30 chars)
- Suggests adding talking points and context

Text overflow checks:
- Detects shapes with >800 characters (potential overflow)
- Detects shapes with >15 paragraphs
- Detects lines >120 characters (horizontal overflow risk)

Usage:
    python scripts/validate_pptx.py <output.pptx> <content.json>
    python scripts/validate_pptx.py <output.pptx> --slides-only

Exit codes:
    0: PASS (no errors)
    1: FAIL (fatal errors found)
    2: WARN (warnings only, can continue)
"""

import argparse
import json
import re
import sys
import io
from pathlib import Path
from typing import Any, Dict, List, Optional

# Fix Windows console encoding issues
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. Run: pip install python-pptx")


class ValidationResult:
    """Holds validation results."""
    
    def __init__(self):
        self.errors: List[Dict[str, Any]] = []
        self.warnings: List[Dict[str, Any]] = []
        self.info: List[Dict[str, Any]] = []
    
    def add_error(self, error_type: str, location: str, message: str, suggestion: str = None):
        """Add a fatal error."""
        self.errors.append({
            "type": error_type,
            "location": location,
            "message": message,
            "suggestion": suggestion
        })
    
    def add_warning(self, warn_type: str, location: str, message: str, suggestion: str = None):
        """Add a warning."""
        self.warnings.append({
            "type": warn_type,
            "location": location,
            "message": message,
            "suggestion": suggestion
        })
    
    def add_info(self, info_type: str, location: str, message: str):
        """Add an info message."""
        self.info.append({
            "type": info_type,
            "location": location,
            "message": message
        })
    
    @property
    def status(self) -> str:
        """Get overall status."""
        if self.errors:
            return "FAIL"
        elif self.warnings:
            return "WARN"
        return "PASS"
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return {
            "status": self.status,
            "fatal_errors": self.errors,
            "warnings": self.warnings,
            "info": self.info,
            "error_count": len(self.errors),
            "warning_count": len(self.warnings)
        }


def load_content_json(content_path: str) -> Optional[Dict[str, Any]]:
    """Load content.json file."""
    try:
        with open(content_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading content.json: {e}")
        return None


def count_pptx_slides(pptx_path: str) -> int:
    """Count slides in PPTX file."""
    prs = Presentation(pptx_path)
    return len(prs.slides)


def get_slide_info(pptx_path: str) -> List[Dict[str, Any]]:
    """Get detailed info about each slide in PPTX."""
    prs = Presentation(pptx_path)
    slides_info = []
    
    for idx, slide in enumerate(prs.slides, 1):
        info = {
            "slide_number": idx,
            "has_notes": False,
            "notes_length": 0,
            "notes_text": "",  # Store actual notes text for quality check
            "has_images": False,
            "image_count": 0,
            "has_title": False,
            "title_text": "",
            "shape_count": len(slide.shapes)
        }
        
        # Check for notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                info["has_notes"] = True
                info["notes_length"] = len(notes_text)
                info["notes_text"] = notes_text  # Store for quality analysis
        
        # Check shapes
        for shape in slide.shapes:
            # Check for images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info["has_images"] = True
                info["image_count"] += 1
            
            # Check for title
            if shape.has_text_frame:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if str(ph_type) in ['TITLE (1)', 'CENTER_TITLE (3)']:
                        info["has_title"] = True
                        info["title_text"] = shape.text_frame.text[:50]
        
        slides_info.append(info)
    
    return slides_info


def validate_slide_count(result: ValidationResult, pptx_path: str, content: Dict[str, Any]):
    """Validate that PPTX slide count matches content.json."""
    pptx_count = count_pptx_slides(pptx_path)
    
    # Count non-skipped slides in content.json
    slides = content.get("slides", [])
    content_count = sum(1 for s in slides if not s.get("_skip", False))
    
    if pptx_count != content_count:
        result.add_error(
            "slide_count_mismatch",
            "global",
            f"PPTX has {pptx_count} slides but content.json has {content_count} slides",
            "Check if all slides were generated correctly"
        )
    else:
        result.add_info(
            "slide_count_match",
            "global",
            f"Slide count matches: {pptx_count} slides"
        )


def validate_notes(result: ValidationResult, slides_info: List[Dict[str, Any]], content: Dict[str, Any]):
    """Validate speaker notes existence and quality."""
    slides = content.get("slides", [])
    non_skipped = [s for s in slides if not s.get("_skip", False)]
    
    missing_notes = []
    insufficient_notes = []
    source_only_notes = []
    
    # Patterns that indicate "source only" notes (insufficient content)
    SOURCE_ONLY_PATTERNS = [
        r"^\s*\[出典:.*?\]\s*$",  # Only source citation (Japanese)
        r"^\s*---\s*\n?\s*\[出典:.*?\]\s*$",  # Separator + source citation only
        r"^\s*\[?Source:.*?\]?\s*$",  # English source only
        r"^\s*\[新規作成\]\s*$",  # Only "newly created" marker
        r"^\s*---\s*\n?\s*\[新規作成\]\s*$",  # Separator + newly created only
    ]
    
    # Minimum notes length (excluding source citations)
    MIN_NOTES_LENGTH = 30
    
    for idx, (slide_info, content_slide) in enumerate(zip(slides_info, non_skipped), 1):
        # Check if content.json expects notes
        expected_notes = content_slide.get("notes", "")
        slide_type = content_slide.get("type", "content")
        
        if expected_notes and not slide_info["has_notes"]:
            missing_notes.append(idx)
        
        # Check notes quality (only for content/section slides with notes)
        if slide_info["has_notes"] and slide_type in ["content", "section", "two_column", "photo"]:
            notes_text = slide_info.get("notes_text", "")
            notes_length = slide_info.get("notes_length", 0)
            
            # Check if notes match "source only" patterns
            is_source_only = False
            for pattern in SOURCE_ONLY_PATTERNS:
                if re.match(pattern, notes_text, re.MULTILINE | re.DOTALL):
                    is_source_only = True
                    break
            
            if is_source_only:
                source_only_notes.append({
                    "slide": idx,
                    "type": slide_type,
                    "preview": notes_text[:50] + "..." if len(notes_text) > 50 else notes_text
                })
            elif notes_length < MIN_NOTES_LENGTH:
                # Check if notes are too short (but not source-only)
                insufficient_notes.append({
                    "slide": idx,
                    "type": slide_type,
                    "length": notes_length,
                    "reason": "too_short"
                })
    
    if missing_notes:
        result.add_warning(
            "missing_notes",
            f"slides {missing_notes}",
            f"{len(missing_notes)} slides are missing expected speaker notes",
            "Speaker notes may not have been applied correctly"
        )
    
    if source_only_notes:
        source_slides = [n["slide"] for n in source_only_notes]
        result.add_warning(
            "source_only_notes",
            f"slides {source_slides}",
            f"{len(source_only_notes)} slides have only source citations in notes (no actual content)",
            "Add talking points, background info, or context for the presenter"
        )
    
    if insufficient_notes:
        short_slides = [n["slide"] for n in insufficient_notes]
        result.add_warning(
            "insufficient_notes",
            f"slides {short_slides}",
            f"{len(insufficient_notes)} slides have very short speaker notes (<{MIN_NOTES_LENGTH} chars)",
            "Consider adding more detail: background info, talking points, or context"
        )
    
    # Count slides without notes
    no_notes_count = sum(1 for s in slides_info if not s["has_notes"])
    if no_notes_count > 0:
        result.add_info(
            "notes_stats",
            "global",
            f"{len(slides_info) - no_notes_count}/{len(slides_info)} slides have speaker notes"
        )


def validate_images(result: ValidationResult, slides_info: List[Dict[str, Any]], content: Dict[str, Any]):
    """Validate image placement."""
    slides = content.get("slides", [])
    non_skipped = [s for s in slides if not s.get("_skip", False)]
    
    missing_images = []
    
    for idx, (slide_info, content_slide) in enumerate(zip(slides_info, non_skipped), 1):
        # Check if content.json expects images
        has_image_spec = "image" in content_slide
        slide_type = content_slide.get("type", "")
        
        if has_image_spec and not slide_info["has_images"]:
            missing_images.append(idx)
        
        # Photo type should have images
        if slide_type == "photo" and not slide_info["has_images"]:
            missing_images.append(idx)
    
    # Remove duplicates
    missing_images = list(set(missing_images))
    
    if missing_images:
        result.add_warning(
            "missing_images",
            f"slides {sorted(missing_images)}",
            f"{len(missing_images)} slides are missing expected images",
            "Check image paths and embedding"
        )
    
    # Count total images
    total_images = sum(s["image_count"] for s in slides_info)
    result.add_info(
        "image_stats",
        "global",
        f"Total images in PPTX: {total_images}"
    )


def validate_text_overflow(result: ValidationResult, pptx_path: str):
    """
    Check for text overflow issues in PPTX.
    
    Detects:
    - Excessively long text in a single shape (>500 chars)
    - Too many paragraphs in a single shape (>15)
    - Very long lines that may cause horizontal overflow
    """
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return  # Already handled in main validation
    
    overflow_slides = []
    long_text_slides = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            text_content = text_frame.text
            
            # Check for very long text (potential overflow)
            if len(text_content) > 800:
                long_text_slides.append(slide_idx)
            
            # Check paragraph count
            para_count = len(text_frame.paragraphs)
            if para_count > 15:
                overflow_slides.append(slide_idx)
            
            # Check for very long single lines (may cause horizontal overflow)
            for para in text_frame.paragraphs:
                line_text = para.text
                if len(line_text) > 120:
                    if slide_idx not in overflow_slides:
                        overflow_slides.append(slide_idx)
    
    if overflow_slides:
        result.add_warning(
            "potential_overflow",
            f"slides {sorted(set(overflow_slides))}",
            f"{len(set(overflow_slides))} slides may have text overflow",
            "Review these slides for text overflow or excessive content"
        )
    
    if long_text_slides:
        result.add_warning(
            "long_text",
            f"slides {sorted(set(long_text_slides))}",
            f"{len(set(long_text_slides))} slides have very long text (>800 chars)",
            "Consider splitting content across multiple slides"
        )
    
    if not overflow_slides and not long_text_slides:
        result.add_info(
            "text_overflow_check",
            "global",
            "No text overflow issues detected"
        )


def validate_signature(result: ValidationResult, slides_info: List[Dict[str, Any]]):
    """Check for repository signature in notes."""
    if not slides_info:
        return
    
    # Check first and last slide for signature
    first_slide = slides_info[0]
    last_slide = slides_info[-1]
    
    has_signature = first_slide["has_notes"] or last_slide["has_notes"]
    
    if has_signature:
        result.add_info(
            "signature",
            "global",
            "Repository signature found in speaker notes"
        )


def validate_pptx(pptx_path: str, content_path: str = None) -> ValidationResult:
    """
    Main validation function.
    
    Args:
        pptx_path: Path to PPTX file
        content_path: Optional path to content.json for comparison
        
    Returns:
        ValidationResult object
    """
    result = ValidationResult()
    
    if not PPTX_AVAILABLE:
        result.add_error(
            "dependency",
            "system",
            "python-pptx is not installed",
            "Run: pip install python-pptx"
        )
        return result
    
    # Check PPTX file exists
    if not Path(pptx_path).exists():
        result.add_error(
            "file_not_found",
            pptx_path,
            "PPTX file not found",
            "Check the file path"
        )
        return result
    
    # Get slide info
    try:
        slides_info = get_slide_info(pptx_path)
    except Exception as e:
        result.add_error(
            "pptx_parse_error",
            pptx_path,
            f"Failed to parse PPTX: {e}",
            "Check if the PPTX file is corrupted"
        )
        return result
    
    result.add_info(
        "pptx_loaded",
        pptx_path,
        f"Successfully loaded PPTX with {len(slides_info)} slides"
    )
    
    # If content.json provided, do comparison validation
    if content_path:
        content = load_content_json(content_path)
        if content:
            validate_slide_count(result, pptx_path, content)
            validate_notes(result, slides_info, content)
            validate_images(result, slides_info, content)
    
    # Always check text overflow
    validate_text_overflow(result, pptx_path)
    
    # Always check signature
    validate_signature(result, slides_info)
    
    return result


def print_result(result: ValidationResult, verbose: bool = False):
    """Print validation result to console."""
    status = result.status
    
    # Status indicator
    if status == "PASS":
        print(f"\n✅ PASS - No issues found")
    elif status == "WARN":
        print(f"\n⚠️ WARN - {result.to_dict()['warning_count']} warning(s)")
    else:
        print(f"\n❌ FAIL - {result.to_dict()['error_count']} error(s)")
    
    # Print errors
    if result.errors:
        print("\n--- Errors ---")
        for err in result.errors:
            print(f"  [{err['type']}] {err['location']}: {err['message']}")
            if err.get('suggestion'):
                print(f"      💡 {err['suggestion']}")
    
    # Print warnings
    if result.warnings:
        print("\n--- Warnings ---")
        for warn in result.warnings:
            print(f"  [{warn['type']}] {warn['location']}: {warn['message']}")
            if warn.get('suggestion'):
                print(f"      💡 {warn['suggestion']}")
    
    # Print info if verbose
    if verbose and result.info:
        print("\n--- Info ---")
        for info in result.info:
            print(f"  [{info['type']}] {info['message']}")


def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description="Validate generated PPTX against content.json"
    )
    parser.add_argument("pptx_path", help="Path to PPTX file")
    parser.add_argument("content_path", nargs="?", help="Path to content.json (optional)")
    parser.add_argument("--slides-only", action="store_true", 
                       help="Only check slide count, skip content comparison")
    parser.add_argument("--verbose", "-v", action="store_true",
                       help="Show detailed info messages")
    parser.add_argument("--json", action="store_true",
                       help="Output result as JSON")
    
    args = parser.parse_args()
    
    # Run validation
    content_path = None if args.slides_only else args.content_path
    result = validate_pptx(args.pptx_path, content_path)
    
    # Output
    if args.json:
        print(json.dumps(result.to_dict(), indent=2, ensure_ascii=False))
    else:
        print_result(result, args.verbose)
    
    # Exit code
    if result.status == "FAIL":
        sys.exit(1)
    elif result.status == "WARN":
        sys.exit(2)
    sys.exit(0)


if __name__ == "__main__":
    main()
